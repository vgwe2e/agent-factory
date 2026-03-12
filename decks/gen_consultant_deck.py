#!/usr/bin/env python3
"""
Agent Factory Consultant Deck — PowerPoint generation script.

Produces a 20-slide brand-compliant Accenture deck showcasing the
agent-factory project as a case study in agentic development.

Output: ./agent-factory-consultant-deck.pptx
"""

import sys, os

# ── Skill imports ──────────────────────────────────────────────
def _find_skill_scripts():
    skills_dir = os.path.join(os.path.expanduser("~"), ".claude", "skills")
    for candidate in ["acnpptx", "pptx"]:
        path = os.path.join(skills_dir, candidate, "scripts")
        if os.path.exists(os.path.join(path, "helpers.py")):
            return path
    raise FileNotFoundError(
        "Cannot find acnpptx skill scripts. "
        f"Checked: {', '.join(os.path.join(skills_dir, n, 'scripts') for n in ['acnpptx', 'pptx'])}"
    )

sys.path.insert(0, _find_skill_scripts())
from helpers import *
from svg_pipeline import add_svg_native as _add_svg_native, cleanup_temp
from pptx import Presentation
from pptx.oxml.ns import qn as _qn
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation(TEMPLATE_PATH)
# Remove pre-existing slides from template
while len(prs.slides) > 0:
    sldId = prs.slides._sldIdLst[0]
    rId = sldId.get(_qn('r:id'))
    prs.part.drop_rel(rId)
    del prs.slides._sldIdLst[0]

reset_slide_counter()

def add_svg_native(slide, svg_str, x, y, w, h, png_width=800):
    return _add_svg_native(slide, prs, svg_str, x, y, w, h, png_width)


# ══════════════════════════════════════════════════════════════
# SVG CHART FUNCTIONS
# ══════════════════════════════════════════════════════════════

def svg_icon_document(size=64, color="#7E00FF"):
    return f'''<svg xmlns="http://www.w3.org/2000/svg" width="{size}" height="{size}" viewBox="0 0 64 64">
  <path d="M16 4h22l14 14v38a4 4 0 0 1-4 4H16a4 4 0 0 1-4-4V8a4 4 0 0 1 4-4z" fill="none" stroke="{color}" stroke-width="2"/>
  <path d="M38 4v14h14" fill="none" stroke="{color}" stroke-width="2" stroke-linejoin="round"/>
  <line x1="20" y1="28" x2="44" y2="28" stroke="{color}" stroke-width="1.5" opacity="0.6"/>
  <line x1="20" y1="36" x2="44" y2="36" stroke="{color}" stroke-width="1.5" opacity="0.6"/>
  <line x1="20" y1="44" x2="36" y2="44" stroke="{color}" stroke-width="1.5" opacity="0.6"/>
</svg>'''

def svg_icon_gear(size=64, color="#7E00FF"):
    return f'''<svg xmlns="http://www.w3.org/2000/svg" width="{size}" height="{size}" viewBox="0 0 64 64">
  <path d="M32 20a12 12 0 1 0 0 24 12 12 0 0 0 0-24zm0 18a6 6 0 1 1 0-12 6 6 0 0 1 0 12z" fill="{color}"/>
  <path d="M34.4 8h-4.8l-1.2 5.2c-1.6.4-3 1.2-4.2 2.2L19.4 13l-3.4 3.4 2.4 4.8c-1 1.2-1.8 2.6-2.2 4.2L11 26.6v4.8l5.2 1.2c.4 1.6 1.2 3 2.2 4.2L16 41.6l3.4 3.4 4.8-2.4c1.2 1 2.6 1.8 4.2 2.2L29.6 50h4.8l1.2-5.2c1.6-.4 3-1.2 4.2-2.2l4.8 2.4 3.4-3.4-2.4-4.8c1-1.2 1.8-2.6 2.2-4.2L53 31.4v-4.8l-5.2-1.2c-.4-1.6-1.2-3-2.2-4.2l2.4-4.8-3.4-3.4-4.8 2.4c-1.2-1-2.6-1.8-4.2-2.2L34.4 8z" fill="none" stroke="{color}" stroke-width="1.5"/>
</svg>'''

def svg_icon_check(size=64, color="#7E00FF"):
    return f'''<svg xmlns="http://www.w3.org/2000/svg" width="{size}" height="{size}" viewBox="0 0 64 64">
  <circle cx="32" cy="32" r="24" fill="none" stroke="{color}" stroke-width="2.5"/>
  <polyline points="20,32 28,42 44,22" fill="none" stroke="{color}" stroke-width="3" stroke-linecap="round" stroke-linejoin="round"/>
</svg>'''

def svg_icon_shield(size=64, color="#7E00FF"):
    return f'''<svg xmlns="http://www.w3.org/2000/svg" width="{size}" height="{size}" viewBox="0 0 64 64">
  <path d="M32 6L8 18v16c0 14 10 22 24 26 14-4 24-12 24-26V18L32 6z" fill="none" stroke="{color}" stroke-width="2.5" stroke-linejoin="round"/>
  <polyline points="22,32 30,40 44,24" fill="none" stroke="{color}" stroke-width="2.5" stroke-linecap="round" stroke-linejoin="round"/>
</svg>'''

def svg_icon_rocket(size=64, color="#7E00FF"):
    return f'''<svg xmlns="http://www.w3.org/2000/svg" width="{size}" height="{size}" viewBox="0 0 64 64">
  <path d="M32 4c-6 8-10 18-10 28 0 4 .5 7 1 10h18c.5-3 1-6 1-10 0-10-4-20-10-28z" fill="none" stroke="{color}" stroke-width="2" stroke-linejoin="round"/>
  <circle cx="32" cy="26" r="4" fill="{color}" opacity="0.6"/>
  <path d="M22 32c-6 2-10 6-10 10h10" fill="none" stroke="{color}" stroke-width="1.5"/>
  <path d="M42 32c6 2 10 6 10 10H42" fill="none" stroke="{color}" stroke-width="1.5"/>
  <path d="M26 52c2 4 4 8 6 8s4-4 6-8" fill="none" stroke="{color}" stroke-width="1.5" opacity="0.7"/>
</svg>'''

def svg_icon_brain(size=64, color="#7E00FF"):
    return f'''<svg xmlns="http://www.w3.org/2000/svg" width="{size}" height="{size}" viewBox="0 0 64 64">
  <path d="M32 8v48" stroke="{color}" stroke-width="1.5" stroke-linecap="round" opacity="0.4"/>
  <path d="M32 12c-4-4-10-4-13 0s-4 10 0 14" fill="none" stroke="{color}" stroke-width="2" stroke-linecap="round" stroke-linejoin="round"/>
  <path d="M19 26c-5 2-8 8-5 13" fill="none" stroke="{color}" stroke-width="2" stroke-linecap="round" stroke-linejoin="round"/>
  <path d="M14 39c-2 5 1 11 7 13h11" fill="none" stroke="{color}" stroke-width="2" stroke-linecap="round" stroke-linejoin="round"/>
  <path d="M32 12c4-4 10-4 13 0s4 10 0 14" fill="none" stroke="{color}" stroke-width="2" stroke-linecap="round" stroke-linejoin="round"/>
  <path d="M45 26c5 2 8 8 5 13" fill="none" stroke="{color}" stroke-width="2" stroke-linecap="round" stroke-linejoin="round"/>
  <path d="M50 39c2 5-1 11-7 13H32" fill="none" stroke="{color}" stroke-width="2" stroke-linecap="round" stroke-linejoin="round"/>
</svg>'''

def svg_icon_code(size=64, color="#7E00FF"):
    return f'''<svg xmlns="http://www.w3.org/2000/svg" width="{size}" height="{size}" viewBox="0 0 64 64">
  <polyline points="20,16 6,32 20,48" fill="none" stroke="{color}" stroke-width="3" stroke-linecap="round" stroke-linejoin="round"/>
  <polyline points="44,16 58,32 44,48" fill="none" stroke="{color}" stroke-width="3" stroke-linecap="round" stroke-linejoin="round"/>
  <line x1="36" y1="10" x2="28" y2="54" stroke="{color}" stroke-width="2.5" stroke-linecap="round" opacity="0.7"/>
</svg>'''

def svg_icon_layers(size=64, color="#7E00FF"):
    return f'''<svg xmlns="http://www.w3.org/2000/svg" width="{size}" height="{size}" viewBox="0 0 64 64">
  <polygon points="32,8 56,22 32,36 8,22" fill="none" stroke="{color}" stroke-width="2" stroke-linejoin="round"/>
  <polygon points="32,20 56,34 32,48 8,34" fill="none" stroke="{color}" stroke-width="2" stroke-linejoin="round" opacity="0.7"/>
  <polygon points="32,32 56,46 32,60 8,46" fill="none" stroke="{color}" stroke-width="2" stroke-linejoin="round" opacity="0.4"/>
</svg>'''

def svg_icon_target(size=64, color="#7E00FF"):
    return f'''<svg xmlns="http://www.w3.org/2000/svg" width="{size}" height="{size}" viewBox="0 0 64 64">
  <circle cx="32" cy="32" r="26" fill="none" stroke="{color}" stroke-width="2"/>
  <circle cx="32" cy="32" r="17" fill="none" stroke="{color}" stroke-width="2" opacity="0.7"/>
  <circle cx="32" cy="32" r="8" fill="none" stroke="{color}" stroke-width="2" opacity="0.5"/>
  <circle cx="32" cy="32" r="2.5" fill="{color}"/>
</svg>'''

def svg_icon_chart_up(size=64, color="#7E00FF"):
    return f'''<svg xmlns="http://www.w3.org/2000/svg" width="{size}" height="{size}" viewBox="0 0 64 64">
  <line x1="8" y1="56" x2="56" y2="56" stroke="#D8D8D8" stroke-width="1.5"/>
  <line x1="8" y1="56" x2="8" y2="8" stroke="#D8D8D8" stroke-width="1.5"/>
  <polyline points="12,48 24,38 36,26 48,12" fill="none" stroke="{color}" stroke-width="2.5" stroke-linecap="round" stroke-linejoin="round"/>
  <circle cx="12" cy="48" r="3" fill="{color}"/>
  <circle cx="24" cy="38" r="3" fill="{color}"/>
  <circle cx="36" cy="26" r="3" fill="{color}"/>
  <circle cx="48" cy="12" r="3" fill="{color}"/>
</svg>'''

def svg_icon_cloud(size=64, color="#7E00FF"):
    return f'''<svg xmlns="http://www.w3.org/2000/svg" width="{size}" height="{size}" viewBox="0 0 64 64">
  <path d="M18 48h30a12 12 0 0 0 0-24 12 12 0 0 0-11.5 8.5A16 16 0 1 0 18 48z" fill="none" stroke="{color}" stroke-width="2.5" stroke-linejoin="round"/>
</svg>'''

def svg_icon_clock(size=64, color="#7E00FF"):
    return f'''<svg xmlns="http://www.w3.org/2000/svg" width="{size}" height="{size}" viewBox="0 0 64 64">
  <circle cx="32" cy="32" r="26" fill="none" stroke="{color}" stroke-width="2.5"/>
  <line x1="32" y1="32" x2="32" y2="16" stroke="{color}" stroke-width="2.5" stroke-linecap="round"/>
  <line x1="32" y1="32" x2="44" y2="38" stroke="{color}" stroke-width="2" stroke-linecap="round"/>
  <circle cx="32" cy="32" r="2" fill="{color}"/>
</svg>'''

def svg_icon_refresh(size=64, color="#7E00FF"):
    return f'''<svg xmlns="http://www.w3.org/2000/svg" width="{size}" height="{size}" viewBox="0 0 64 64">
  <path d="M48 16A22 22 0 0 0 12 26" fill="none" stroke="{color}" stroke-width="2.5" stroke-linecap="round"/>
  <path d="M16 48A22 22 0 0 0 52 38" fill="none" stroke="{color}" stroke-width="2.5" stroke-linecap="round"/>
  <polyline points="12,16 12,28 24,28" fill="none" stroke="{color}" stroke-width="2.5" stroke-linecap="round" stroke-linejoin="round"/>
  <polyline points="52,48 52,36 40,36" fill="none" stroke="{color}" stroke-width="2.5" stroke-linecap="round" stroke-linejoin="round"/>
</svg>'''

def svg_icon_lightbulb(size=64, color="#7E00FF"):
    return f'''<svg xmlns="http://www.w3.org/2000/svg" width="{size}" height="{size}" viewBox="0 0 64 64">
  <path d="M24 40c-1.5-3-6-7-6-16a14 14 0 1 1 28 0c0 9-4.5 13-6 16z" fill="none" stroke="{color}" stroke-width="2" stroke-linejoin="round"/>
  <line x1="24" y1="40" x2="40" y2="40" stroke="{color}" stroke-width="2"/>
  <line x1="25" y1="46" x2="39" y2="46" stroke="{color}" stroke-width="2" stroke-linecap="round"/>
  <line x1="27" y1="52" x2="37" y2="52" stroke="{color}" stroke-width="2" stroke-linecap="round"/>
</svg>'''

def svg_icon_handshake(size=64, color="#7E00FF"):
    return f'''<svg xmlns="http://www.w3.org/2000/svg" width="{size}" height="{size}" viewBox="0 0 64 64">
  <path d="M4 24h10l8-4h6l8 8 8-4h6l10 4v4H50l-4 4-6 6-4 2h-4l-6-4-6-2H4V24z" fill="none" stroke="{color}" stroke-width="2" stroke-linejoin="round" stroke-linecap="round"/>
  <path d="M22 20l10 2 12-2" fill="none" stroke="{color}" stroke-width="1.5" stroke-linecap="round" stroke-linejoin="round"/>
  <line x1="4" y1="24" x2="4" y2="34" stroke="{color}" stroke-width="2"/>
  <line x1="60" y1="24" x2="60" y2="34" stroke="{color}" stroke-width="2"/>
  <path d="M26 38l6 6" fill="none" stroke="{color}" stroke-width="1.5" stroke-linecap="round"/>
  <path d="M32 34l6 6" fill="none" stroke="{color}" stroke-width="1.5" stroke-linecap="round"/>
</svg>'''

def svg_icon_people(size=64, color="#7E00FF"):
    return f'''<svg xmlns="http://www.w3.org/2000/svg" width="{size}" height="{size}" viewBox="0 0 64 64">
  <circle cx="22" cy="18" r="7" fill="none" stroke="{color}" stroke-width="2"/>
  <path d="M8 48c0-8 6-14 14-14s14 6 14 14" fill="none" stroke="{color}" stroke-width="2" stroke-linecap="round"/>
  <circle cx="42" cy="18" r="7" fill="none" stroke="{color}" stroke-width="2"/>
  <path d="M28 48c0-8 6-14 14-14s14 6 14 14" fill="none" stroke="{color}" stroke-width="2" stroke-linecap="round"/>
</svg>'''

def svg_icon_database(size=64, color="#7E00FF"):
    return f'''<svg xmlns="http://www.w3.org/2000/svg" width="{size}" height="{size}" viewBox="0 0 64 64">
  <ellipse cx="32" cy="14" rx="22" ry="8" fill="none" stroke="{color}" stroke-width="2"/>
  <path d="M10 14v36c0 4.4 9.8 8 22 8s22-3.6 22-8V14" fill="none" stroke="{color}" stroke-width="2"/>
  <path d="M10 26c0 4.4 9.8 8 22 8s22-3.6 22-8" fill="none" stroke="{color}" stroke-width="1.5" opacity="0.6"/>
  <path d="M10 38c0 4.4 9.8 8 22 8s22-3.6 22-8" fill="none" stroke="{color}" stroke-width="1.5" opacity="0.6"/>
</svg>'''

def svg_icon_network(size=64, color="#7E00FF"):
    return f'''<svg xmlns="http://www.w3.org/2000/svg" width="{size}" height="{size}" viewBox="0 0 64 64">
  <circle cx="32" cy="14" r="7" fill="none" stroke="{color}" stroke-width="2.5"/>
  <circle cx="14" cy="50" r="7" fill="none" stroke="{color}" stroke-width="2.5"/>
  <circle cx="50" cy="50" r="7" fill="none" stroke="{color}" stroke-width="2.5"/>
  <circle cx="32" cy="14" r="3" fill="{color}"/>
  <circle cx="14" cy="50" r="3" fill="{color}"/>
  <circle cx="50" cy="50" r="3" fill="{color}"/>
  <line x1="32" y1="21" x2="14" y2="43" stroke="{color}" stroke-width="1.5"/>
  <line x1="32" y1="21" x2="50" y2="43" stroke="{color}" stroke-width="1.5"/>
  <line x1="21" y1="50" x2="43" y2="50" stroke="{color}" stroke-width="1.5"/>
</svg>'''

def svg_icon_search(size=64, color="#7E00FF"):
    return f'''<svg xmlns="http://www.w3.org/2000/svg" width="{size}" height="{size}" viewBox="0 0 64 64">
  <circle cx="28" cy="28" r="18" fill="none" stroke="{color}" stroke-width="2.5"/>
  <line x1="40.7" y1="40.7" x2="56" y2="56" stroke="{color}" stroke-width="3" stroke-linecap="round"/>
</svg>'''

def svg_icon_calendar(size=64, color="#7E00FF"):
    return f'''<svg xmlns="http://www.w3.org/2000/svg" width="{size}" height="{size}" viewBox="0 0 64 64">
  <rect x="8" y="14" width="48" height="42" rx="4" fill="none" stroke="{color}" stroke-width="2"/>
  <line x1="8" y1="26" x2="56" y2="26" stroke="{color}" stroke-width="2"/>
  <line x1="20" y1="8" x2="20" y2="20" stroke="{color}" stroke-width="2.5" stroke-linecap="round"/>
  <line x1="44" y1="8" x2="44" y2="20" stroke="{color}" stroke-width="2.5" stroke-linecap="round"/>
  <circle cx="20" cy="35" r="2" fill="{color}"/>
  <circle cx="32" cy="35" r="2" fill="{color}"/>
  <circle cx="44" cy="35" r="2" fill="{color}"/>
  <circle cx="20" cy="46" r="2" fill="{color}"/>
  <circle cx="32" cy="46" r="2" fill="{color}"/>
  <circle cx="44" cy="46" r="2" fill="{color}"/>
</svg>'''

# SVG chart functions
def svg_chevron_flow(items, w=800, h=70, color="#7E00FF"):
    n = len(items)
    cw = w / n
    arrow = 14
    shapes = ""
    for i, lbl in enumerate(items):
        x = i * cw
        op = max(0.5, 1.0 - i * 0.12)
        shapes += (f'<polygon points="{x},0 {x + cw - arrow},0 '
                   f'{x + cw},{h / 2} {x + cw - arrow},{h} '
                   f'{x},{h}" fill="{color}" opacity="{op}"/>')
        shapes += (f'<text x="{x + cw / 2}" '
                   f'y="{h / 2 + 5}" text-anchor="middle" '
                   f'font-family="Graphik, sans-serif" font-size="13" '
                   f'fill="white" font-weight="bold">{lbl}</text>')
    return (f'<svg xmlns="http://www.w3.org/2000/svg" '
            f'width="{w}" height="{h}">{shapes}</svg>')

def svg_layers_diagram(layers, w=500, h=240, color="#7E00FF"):
    n = len(layers)
    lh = (h - (n - 1) * 6) / n
    shapes = ""
    for i, lbl in enumerate(layers):
        y = i * (lh + 6)
        op = max(0.35, 1.0 - i * 0.18)
        shapes += (f'<rect x="20" y="{y}" width="{w - 40}" height="{lh}" '
                   f'rx="6" fill="{color}" opacity="{op}"/>')
        shapes += (f'<text x="{w / 2}" y="{y + lh / 2 + 5}" '
                   f'text-anchor="middle" font-family="Graphik, sans-serif" '
                   f'font-size="14" fill="white" font-weight="bold">{lbl}</text>')
    return (f'<svg xmlns="http://www.w3.org/2000/svg" '
            f'width="{w}" height="{int(n * (lh + 6) - 6)}">{shapes}</svg>')

def svg_callout_badge(text="1", size=40, bg_color="#5F0095", text_color="#FFFFFF"):
    return f'''<svg xmlns="http://www.w3.org/2000/svg" width="{size}" height="{size}" viewBox="0 0 {size} {size}">
  <circle cx="{size//2}" cy="{size//2}" r="{size//2-2}" fill="{bg_color}"/>
  <text x="{size//2}" y="{size//2+1}" text-anchor="middle" dominant-baseline="central" font-family="Graphik, sans-serif" font-size="{int(size*0.45)}" font-weight="bold" fill="{text_color}">{text}</text>
</svg>'''

def svg_divider_gradient(w=800, color="#7E00FF"):
    return f'''<svg xmlns="http://www.w3.org/2000/svg" width="{w}" height="6" viewBox="0 0 {w} 6">
  <defs>
    <linearGradient id="dg" x1="0%" y1="0%" x2="100%" y2="0%">
      <stop offset="0%" stop-color="{color}" stop-opacity="1"/>
      <stop offset="100%" stop-color="{color}" stop-opacity="0.1"/>
    </linearGradient>
  </defs>
  <rect x="0" y="0" width="{w}" height="6" rx="3" fill="url(#dg)"/>
</svg>'''

def svg_highlight_bar(w=800, h=60, bg_color="#F3E8FF", border_color="#7E00FF", arrow_color="#7E00FF"):
    return f'''<svg xmlns="http://www.w3.org/2000/svg" width="{w}" height="{h}" viewBox="0 0 {w} {h}">
  <rect x="0" y="0" width="{w}" height="{h}" rx="4" fill="{bg_color}" stroke="{border_color}" stroke-width="2"/>
  <polygon points="4,{h//2-8} 20,{h//2} 4,{h//2+8}" fill="{arrow_color}"/>
</svg>'''


def svg_architecture_diagram():
    """Architecture diagram: two systems side by side."""
    w, h = 800, 340
    svg = f'''<svg xmlns="http://www.w3.org/2000/svg" width="{w}" height="{h}" viewBox="0 0 {w} {h}">
  <!-- Left system: Aera Engine -->
  <rect x="20" y="10" width="370" height="320" rx="8" fill="#F5F5F5" stroke="#D8D8D8" stroke-width="1.5"/>
  <rect x="20" y="10" width="370" height="40" rx="8" fill="#5F0095"/>
  <text x="205" y="36" text-anchor="middle" font-family="Graphik, sans-serif" font-size="15" fill="white" font-weight="bold">Aera Skill Feasibility Engine (src/)</text>

  <!-- Pipeline stages -->
  <rect x="40" y="65" width="330" height="35" rx="4" fill="#7E00FF" opacity="0.9"/>
  <text x="205" y="87" text-anchor="middle" font-family="Graphik, sans-serif" font-size="13" fill="white" font-weight="bold">CLI Entry (commander + Zod validation)</text>

  <rect x="40" y="110" width="155" height="35" rx="4" fill="#7E00FF" opacity="0.75"/>
  <text x="117" y="132" text-anchor="middle" font-family="Graphik, sans-serif" font-size="12" fill="white" font-weight="bold">8B Triage</text>

  <rect x="215" y="110" width="155" height="35" rx="4" fill="#7E00FF" opacity="0.75"/>
  <text x="292" y="132" text-anchor="middle" font-family="Graphik, sans-serif" font-size="12" fill="white" font-weight="bold">32B Scoring</text>

  <rect x="40" y="155" width="330" height="35" rx="4" fill="#7E00FF" opacity="0.6"/>
  <text x="205" y="177" text-anchor="middle" font-family="Graphik, sans-serif" font-size="13" fill="white" font-weight="bold">Simulation (Mermaid + YAML + Mock Tests)</text>

  <rect x="40" y="200" width="330" height="35" rx="4" fill="#7E00FF" opacity="0.5"/>
  <text x="205" y="222" text-anchor="middle" font-family="Graphik, sans-serif" font-size="13" fill="white" font-weight="bold">Reports + Git Auto-Commit</text>

  <rect x="40" y="250" width="155" height="65" rx="4" fill="#F3E8FF" stroke="#7E00FF" stroke-width="1"/>
  <text x="117" y="275" text-anchor="middle" font-family="Graphik, sans-serif" font-size="11" fill="#5F0095" font-weight="bold">Ollama (local)</text>
  <text x="117" y="295" text-anchor="middle" font-family="Graphik, sans-serif" font-size="11" fill="#666666">8B + 32B models</text>

  <rect x="215" y="250" width="155" height="65" rx="4" fill="#F3E8FF" stroke="#7E00FF" stroke-width="1"/>
  <text x="292" y="275" text-anchor="middle" font-family="Graphik, sans-serif" font-size="11" fill="#5F0095" font-weight="bold">vLLM (cloud)</text>
  <text x="292" y="295" text-anchor="middle" font-family="Graphik, sans-serif" font-size="11" fill="#666666">H100 via RunPod</text>

  <!-- Right system: Agent Harness -->
  <rect x="410" y="10" width="370" height="320" rx="8" fill="#F5F5F5" stroke="#D8D8D8" stroke-width="1.5"/>
  <rect x="410" y="10" width="370" height="40" rx="8" fill="#5F0095"/>
  <text x="595" y="36" text-anchor="middle" font-family="Graphik, sans-serif" font-size="15" fill="white" font-weight="bold">Agent Harness (seed/)</text>

  <rect x="430" y="65" width="330" height="35" rx="4" fill="#7E00FF" opacity="0.9"/>
  <text x="595" y="87" text-anchor="middle" font-family="Graphik, sans-serif" font-size="13" fill="white" font-weight="bold">SSE API (Next.js POST /api/chat)</text>

  <rect x="430" y="110" width="330" height="35" rx="4" fill="#7E00FF" opacity="0.75"/>
  <text x="595" y="132" text-anchor="middle" font-family="Graphik, sans-serif" font-size="13" fill="white" font-weight="bold">Orchestrator (async generator loop)</text>

  <rect x="430" y="155" width="155" height="35" rx="4" fill="#7E00FF" opacity="0.6"/>
  <text x="507" y="177" text-anchor="middle" font-family="Graphik, sans-serif" font-size="12" fill="white" font-weight="bold">Providers</text>

  <rect x="605" y="155" width="155" height="35" rx="4" fill="#7E00FF" opacity="0.6"/>
  <text x="682" y="177" text-anchor="middle" font-family="Graphik, sans-serif" font-size="12" fill="white" font-weight="bold">Tool Registry</text>

  <rect x="430" y="200" width="330" height="35" rx="4" fill="#7E00FF" opacity="0.5"/>
  <text x="595" y="222" text-anchor="middle" font-family="Graphik, sans-serif" font-size="13" fill="white" font-weight="bold">Config (per-agent settings)</text>

  <text x="595" y="270" text-anchor="middle" font-family="Graphik, sans-serif" font-size="12" fill="#666666">Hand-written orchestrator</text>
  <text x="595" y="290" text-anchor="middle" font-family="Graphik, sans-serif" font-size="12" fill="#666666">No LangChain / CrewAI / Autogen</text>
  <text x="595" y="310" text-anchor="middle" font-family="Graphik, sans-serif" font-size="12" fill="#5F0095" font-weight="bold">~150 LOC async generators</text>
</svg>'''
    return svg


def svg_scoring_comparison():
    """Bar chart comparing adoption-first vs traditional weighting."""
    w, h = 800, 280
    # Bars for "Our approach" and "Traditional"
    bar_w = 70
    gap = 20
    group_w = bar_w * 2 + gap
    labels = ["Technical\nFeasibility", "Adoption\nRealism", "Value &amp;\nEfficiency"]
    ours = [0.30, 0.45, 0.25]
    trad = [0.50, 0.30, 0.20]
    max_val = 0.55
    plot_h = 180
    plot_y = 30
    base_y = plot_y + plot_h
    start_x = 100

    bars = ""
    for i, (lbl, o, t) in enumerate(zip(labels, ours, trad)):
        gx = start_x + i * (group_w + 80)
        # Traditional bar
        th = (t / max_val) * plot_h
        bars += f'<rect x="{gx}" y="{base_y - th}" width="{bar_w}" height="{th}" fill="#D8D8D8" rx="3"/>'
        bars += f'<text x="{gx + bar_w/2}" y="{base_y - th - 8}" text-anchor="middle" font-family="Graphik, sans-serif" font-size="13" fill="#666666" font-weight="bold">{int(t*100)}%</text>'
        # Our bar
        oh = (o / max_val) * plot_h
        bars += f'<rect x="{gx + bar_w + gap}" y="{base_y - oh}" width="{bar_w}" height="{oh}" fill="#7E00FF" rx="3"/>'
        bars += f'<text x="{gx + bar_w + gap + bar_w/2}" y="{base_y - oh - 8}" text-anchor="middle" font-family="Graphik, sans-serif" font-size="13" fill="#5F0095" font-weight="bold">{int(o*100)}%</text>'
        # Label
        lines = lbl.split("\n")
        for li, line in enumerate(lines):
            bars += f'<text x="{gx + group_w/2}" y="{base_y + 22 + li * 16}" text-anchor="middle" font-family="Graphik, sans-serif" font-size="12" fill="#333333">{line}</text>'

    # Axis line
    bars += f'<line x1="{start_x - 20}" y1="{base_y}" x2="{start_x + 3 * (group_w + 80) - 60}" y2="{base_y}" stroke="#D8D8D8" stroke-width="1"/>'

    # Legend
    bars += f'<rect x="{w - 250}" y="10" width="16" height="16" fill="#D8D8D8" rx="2"/>'
    bars += f'<text x="{w - 228}" y="23" font-family="Graphik, sans-serif" font-size="12" fill="#666666">Traditional (tech-first)</text>'
    bars += f'<rect x="{w - 250}" y="34" width="16" height="16" fill="#7E00FF" rx="2"/>'
    bars += f'<text x="{w - 228}" y="47" font-family="Graphik, sans-serif" font-size="12" fill="#5F0095">Ours (adoption-first)</text>'

    return f'<svg xmlns="http://www.w3.org/2000/svg" width="{w}" height="{h}" viewBox="0 0 {w} {h}">{bars}</svg>'


def svg_metrics_dashboard():
    """Grid of key metrics for v1.0."""
    w, h = 800, 300
    metrics = [
        ("213K", "Lines of Code", "TypeScript (strict)"),
        ("412", "Tests Passing", "Node.js built-in runner"),
        ("44/44", "Requirements", "Fully satisfied"),
        ("31", "Atomic Plans", "Across 11 phases"),
        ("93", "Minutes Total", "Automated execution"),
        ("3 min", "Avg per Plan", "Including TDD"),
    ]
    cols, rows_count = 3, 2
    cw_px = (w - 40) / cols
    rh = (h - 20) / rows_count
    cells = ""
    for i, (val, label, sub) in enumerate(metrics):
        col = i % cols
        row = i // cols
        cx = 20 + col * cw_px
        cy = 10 + row * rh
        cells += f'<rect x="{cx + 4}" y="{cy + 4}" width="{cw_px - 8}" height="{rh - 8}" rx="6" fill="#F5F5F5" stroke="#D8D8D8" stroke-width="1"/>'
        cells += f'<text x="{cx + cw_px/2}" y="{cy + rh * 0.35}" text-anchor="middle" font-family="Graphik, sans-serif" font-size="28" fill="#5F0095" font-weight="bold">{val}</text>'
        cells += f'<text x="{cx + cw_px/2}" y="{cy + rh * 0.58}" text-anchor="middle" font-family="Graphik, sans-serif" font-size="14" fill="#333333" font-weight="bold">{label}</text>'
        cells += f'<text x="{cx + cw_px/2}" y="{cy + rh * 0.78}" text-anchor="middle" font-family="Graphik, sans-serif" font-size="12" fill="#666666">{sub}</text>'
    return f'<svg xmlns="http://www.w3.org/2000/svg" width="{w}" height="{h}" viewBox="0 0 {w} {h}">{cells}</svg>'


def svg_perf_comparison():
    """Side-by-side comparison: local vs cloud scoring."""
    w, h = 800, 260
    # Two bars: Local (17hr) vs Cloud (<30min)
    bar_max_w = 600
    local_w = bar_max_w
    cloud_w = int(bar_max_w * (30 / (17 * 60)))  # proportional
    cloud_w = max(cloud_w, 40)  # min visible width
    bar_h = 50
    y1 = 40
    y2 = 130
    lx = 160

    svg = f'''<svg xmlns="http://www.w3.org/2000/svg" width="{w}" height="{h}" viewBox="0 0 {w} {h}">
  <text x="20" y="{y1 + 32}" font-family="Graphik, sans-serif" font-size="14" fill="#333333" font-weight="bold">Local</text>
  <text x="20" y="{y1 + 52}" font-family="Graphik, sans-serif" font-size="12" fill="#666666">(Ollama)</text>
  <rect x="{lx}" y="{y1}" width="{local_w}" height="{bar_h}" rx="4" fill="#D8D8D8"/>
  <text x="{lx + local_w/2}" y="{y1 + 30}" text-anchor="middle" font-family="Graphik, sans-serif" font-size="16" fill="#333333" font-weight="bold">~17 hours (339 opportunities)</text>

  <text x="20" y="{y2 + 32}" font-family="Graphik, sans-serif" font-size="14" fill="#333333" font-weight="bold">Cloud</text>
  <text x="20" y="{y2 + 52}" font-family="Graphik, sans-serif" font-size="12" fill="#666666">(vLLM H100)</text>
  <rect x="{lx}" y="{y2}" width="{cloud_w}" height="{bar_h}" rx="4" fill="#7E00FF"/>
  <text x="{lx + cloud_w + 12}" y="{y2 + 30}" font-family="Graphik, sans-serif" font-size="16" fill="#5F0095" font-weight="bold">&lt; 30 min (10-20 concurrent)</text>

  <text x="{lx}" y="{h - 20}" font-family="Graphik, sans-serif" font-size="12" fill="#666666" font-style="italic">Cost target: &lt; $10/run on RunPod ephemeral H100 ($2-3/hr)</text>
  <text x="{lx + local_w - 40}" y="{y1 - 10}" font-family="Graphik, sans-serif" font-size="12" fill="#666666">34x faster</text>
</svg>'''
    return svg


# ══════════════════════════════════════════════════════════════
# HELPER: Add text to a textbox with standard formatting
# ══════════════════════════════════════════════════════════════

def _add_body_text(slide, x, y, w, h, lines, font_size=Pt(14), color=TEXT_BODY, bold=False, spacing=Pt(6)):
    """Add a text box with multiple lines."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = spacing if i > 0 else Pt(0)
        run = p.add_run()
        run.text = line
        run.font.size = font_size
        run.font.color.rgb = color
        run.font.bold = bold
        run.font.name = FONT
    return tb


def _add_card(slide, x, y, w, h, title, bullets, icon_svg=None):
    """Add a card with purple top bar, optional icon, title, and bullets."""
    # Card background
    card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h))
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = LIGHT_GRAY
    card.line.width = Pt(0.75)

    # Purple top bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(0.05))
    bar.fill.solid()
    bar.fill.fore_color.rgb = PURPLE
    bar.line.fill.background()

    text_y = y + 0.15
    text_x = x + 0.15

    # Icon
    if icon_svg:
        add_svg_native(slide, icon_svg, x=x + 0.12, y=y + 0.15, w=0.35, h=0.35, png_width=256)
        text_y = y + 0.55

    # Title
    ttb = slide.shapes.add_textbox(Inches(text_x), Inches(text_y), Inches(w - 0.30), Inches(0.35))
    tf = ttb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.size = Pt(14)
    r.font.bold = True
    r.font.color.rgb = DEEP_PURPLE
    r.font.name = FONT

    # Bullets
    bullet_y = text_y + 0.35
    btb = slide.shapes.add_textbox(Inches(text_x), Inches(bullet_y), Inches(w - 0.30), Inches(h - (bullet_y - y) - 0.10))
    bf = btb.text_frame
    bf.word_wrap = True
    for i, b in enumerate(bullets):
        p = bf.paragraphs[0] if i == 0 else bf.add_paragraph()
        p.space_before = Pt(4)
        r = p.add_run()
        r.text = b
        r.font.size = Pt(14)
        r.font.color.rgb = TEXT_BODY
        r.font.name = FONT


# ══════════════════════════════════════════════════════════════
# SLIDE 1: COVER
# ══════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[0])

# Title — BLACK text
for ph in slide.placeholders:
    idx = ph.placeholder_format.idx
    if idx == 0:  # CENTER_TITLE
        ph.text = "Agentic Development: A Case Study"
        for p in ph.text_frame.paragraphs:
            p.font.size = Pt(40)
            p.font.bold = True
            p.font.color.rgb = BLACK
            p.font.name = FONT
    elif idx == 1:  # SUBTITLE
        ph.text = "How Claude Code + fine-grained planning shipped 213K LOC in 93 minutes"
        for p in ph.text_frame.paragraphs:
            p.font.size = Pt(18)
            p.font.color.rgb = BLACK
            p.font.name = FONT
    elif idx == 12:  # BODY
        ph.text = "Autonomous software delivery with test-driven development, checkpoint recovery, and zero frameworks"
        for p in ph.text_frame.paragraphs:
            p.font.size = Pt(14)
            p.font.color.rgb = BLACK
            p.font.name = FONT
    elif idx == 2:  # DATE
        ph.text = "March 2026"
        for p in ph.text_frame.paragraphs:
            p.font.size = Pt(14)
            p.font.color.rgb = TEXT_SUB
            p.font.name = FONT

add_logo_to_cover(slide)


# ══════════════════════════════════════════════════════════════
# SLIDE 2: THE ENTERPRISE AI GAP
# ══════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[2])
for ph in slide.placeholders:
    idx = ph.placeholder_format.idx
    if idx == 11:
        ph.text = "THE PROBLEM"
        for p in ph.text_frame.paragraphs:
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = PURPLE
            p.font.name = FONT
    elif idx == 0:
        ph.text = "Most AI features fail at adoption, not technology"
        for p in ph.text_frame.paragraphs:
            p.font.size = Pt(28)
            p.font.bold = True
            p.font.color.rgb = BLACK
            p.font.name = FONT
    elif idx == 10:
        ph.text = "Technical feasibility scoring alone produces features nobody uses. The gap is between what can be built and what will be adopted."
        for p in ph.text_frame.paragraphs:
            p.font.size = Pt(18)
            p.font.color.rgb = BLACK
            p.font.name = FONT

# Pattern D: Two-column split
col_w = (CW - 0.20) / 2
col_h = AH - 0.60

# Left card: The problem
_add_card(slide, ML, CY + 0.55, col_w, col_h, "Traditional approach",
    ["Technical feasibility weighted 50%",
     "Adoption and change management weighted 30%",
     "Value metrics weighted 20%",
     "Result: technically sound features that real users reject or ignore"],
    svg_icon_search())

# Right card: Our approach
_add_card(slide, ML + col_w + 0.20, CY + 0.55, col_w, col_h, "Adoption-first approach",
    ["Technical feasibility weighted 30%",
     "Adoption realism weighted 45% (highest)",
     "Value and efficiency weighted 25%",
     "Result: features scored on whether organizations will actually use them"],
    svg_icon_target())

set_footer(slide)
add_gt_to_slide(slide)


# ══════════════════════════════════════════════════════════════
# SLIDE 4: WHAT WE SET OUT TO BUILD
# ══════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[2])
for ph in slide.placeholders:
    idx = ph.placeholder_format.idx
    if idx == 11:
        ph.text = "THE PROBLEM"
        for p in ph.text_frame.paragraphs:
            p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = PURPLE; p.font.name = FONT
    elif idx == 0:
        ph.text = "An offline-first engine that scores with adoption realism"
        for p in ph.text_frame.paragraphs:
            p.font.size = Pt(28); p.font.bold = True; p.font.color.rgb = BLACK; p.font.name = FONT
    elif idx == 10:
        ph.text = "Single CLI command evaluates an Aera enterprise hierarchy export across three weighted lenses, runs unattended overnight."
        for p in ph.text_frame.paragraphs:
            p.font.size = Pt(18); p.font.color.rgb = BLACK; p.font.name = FONT

# Pattern I: Chevron process flow
flow_svg = svg_chevron_flow(["Ingestion", "Triage", "Scoring", "Simulation", "Reports"], w=800, h=60)
add_svg_native(slide, flow_svg, x=ML, y=CY + 0.55, w=CW, h=0.55)

# Details below chevrons
details = [
    ("Zod-validated JSON", "Parse any client\nhierarchy export"),
    ("5 Red Flags + 3 Tiers", "Filter bad opps\nbefore LLM scoring"),
    ("3-Lens Composite", "Tech + Adoption\n+ Value scoring"),
    ("4 Artifact Types", "Mermaid, YAML,\nmock tests, specs"),
    ("Git Auto-Commit", "Overnight-safe\nwith checkpoints"),
]
dcw = (CW - 4 * 0.12) / 5
for i, (title, desc) in enumerate(details):
    dx = ML + i * (dcw + 0.12)
    dy = CY + 1.25

    card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(dx), Inches(dy), Inches(dcw), Inches(3.70))
    card.fill.solid()
    card.fill.fore_color.rgb = BG_LIGHT
    card.line.color.rgb = LIGHT_GRAY
    card.line.width = Pt(0.75)

    ttb = slide.shapes.add_textbox(Inches(dx + 0.10), Inches(dy + 0.10), Inches(dcw - 0.20), Inches(0.40))
    tf = ttb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = title; r.font.size = Pt(14); r.font.bold = True; r.font.color.rgb = DEEP_PURPLE; r.font.name = FONT

    dtb = slide.shapes.add_textbox(Inches(dx + 0.10), Inches(dy + 0.55), Inches(dcw - 0.20), Inches(3.00))
    df = dtb.text_frame; df.word_wrap = True
    p = df.paragraphs[0]
    r = p.add_run(); r.text = desc; r.font.size = Pt(14); r.font.color.rgb = TEXT_BODY; r.font.name = FONT

set_footer(slide)
add_gt_to_slide(slide)


# ══════════════════════════════════════════════════════════════
# SLIDE 5: SYSTEM ARCHITECTURE OVERVIEW
# ══════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[2])
for ph in slide.placeholders:
    idx = ph.placeholder_format.idx
    if idx == 11:
        ph.text = "ARCHITECTURE"
        for p in ph.text_frame.paragraphs:
            p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = PURPLE; p.font.name = FONT
    elif idx == 0:
        ph.text = "Two complementary systems in one repository"
        for p in ph.text_frame.paragraphs:
            p.font.size = Pt(28); p.font.bold = True; p.font.color.rgb = BLACK; p.font.name = FONT
    elif idx == 10:
        ph.text = "The Aera Engine (src/) evaluates enterprise opportunities. The Agent Harness (seed/) provides the agentic loop pattern with zero framework dependencies."
        for p in ph.text_frame.paragraphs:
            p.font.size = Pt(18); p.font.color.rgb = BLACK; p.font.name = FONT

arch_svg = svg_architecture_diagram()
add_svg_native(slide, arch_svg, x=ML, y=CY + 0.55, w=CW, h=4.50)

set_footer(slide)
add_gt_to_slide(slide)


# ══════════════════════════════════════════════════════════════
# SLIDE 6: THREE-LENS SCORING
# ══════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[2])
for ph in slide.placeholders:
    idx = ph.placeholder_format.idx
    if idx == 11:
        ph.text = "ARCHITECTURE"
        for p in ph.text_frame.paragraphs:
            p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = PURPLE; p.font.name = FONT
    elif idx == 0:
        ph.text = "Adoption realism is the highest-weighted lens"
        for p in ph.text_frame.paragraphs:
            p.font.size = Pt(28); p.font.bold = True; p.font.color.rgb = BLACK; p.font.name = FONT
    elif idx == 10:
        ph.text = "Weighting adoption at 0.45 prevents the common failure mode: technically sound features that nobody uses. A 0.60 composite threshold gates promotion to simulation."
        for p in ph.text_frame.paragraphs:
            p.font.size = Pt(18); p.font.color.rgb = BLACK; p.font.name = FONT

scoring_svg = svg_scoring_comparison()
add_svg_native(slide, scoring_svg, x=ML, y=CY + 0.55, w=CW, h=3.60)

# Key message bar at bottom
kmb_y = BY - 0.70
bar_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
    Inches(ML), Inches(kmb_y), Inches(CW), Inches(0.65))
bar_bg.fill.solid()
bar_bg.fill.fore_color.rgb = LP_BG
bar_bg.line.color.rgb = PURPLE
bar_bg.line.width = Pt(1)
kmtb = slide.shapes.add_textbox(Inches(ML + 0.20), Inches(kmb_y + 0.10), Inches(CW - 0.40), Inches(0.45))
kf = kmtb.text_frame; kf.word_wrap = True
p = kf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = "Composite = (Technical x 0.30) + (Adoption x 0.45) + (Value x 0.25) -- threshold 0.60 to advance"
r.font.size = Pt(14); r.font.bold = True; r.font.color.rgb = DEEP_PURPLE; r.font.name = FONT

set_footer(slide)
add_gt_to_slide(slide)


# ══════════════════════════════════════════════════════════════
# SLIDE 7: FIVE RED FLAGS (AUTOMATED TRIAGE)
# ══════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[2])
for ph in slide.placeholders:
    idx = ph.placeholder_format.idx
    if idx == 11:
        ph.text = "ARCHITECTURE"
        for p in ph.text_frame.paragraphs:
            p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = PURPLE; p.font.name = FONT
    elif idx == 0:
        ph.text = "Five red flags filter bad opportunities early"
        for p in ph.text_frame.paragraphs:
            p.font.size = Pt(28); p.font.bold = True; p.font.color.rgb = BLACK; p.font.name = FONT
    elif idx == 10:
        ph.text = "Pure-function triage runs before expensive LLM scoring. Bad opportunities are skipped, demoted, or flagged automatically."
        for p in ph.text_frame.paragraphs:
            p.font.size = Pt(18); p.font.color.rgb = BLACK; p.font.name = FONT

# Pattern C: Data table
tbl = slide.shapes.add_table(6, 4, Inches(ML), Inches(CY + 0.55), Inches(CW), Inches(4.40)).table
tbl.columns[0].width = Inches(2.50)
tbl.columns[1].width = Inches(4.00)
tbl.columns[2].width = Inches(2.50)
tbl.columns[3].width = Inches(3.50)

headers = ["Red Flag", "Detection Rule", "Action", "Purpose"]
flags = [
    ("DEAD_ZONE", "0% decision density across all L4s", "Skip", "No decisions = no skill to build"),
    ("PHANTOM", "opportunity_exists = false", "Skip", "Opportunity does not actually exist"),
    ("NO_STAKES", "Zero HIGH financial + all SECOND order", "Demote", "Too low-value for LLM investment"),
    ("CONFIDENCE_GAP", ">50% of L4s have LOW confidence", "Flag", "Data quality too poor for reliable scoring"),
    ("ORPHAN", "l4_count < 3", "Flag", "Too few activities for meaningful analysis"),
]

for ci, h_text in enumerate(headers):
    cell = tbl.cell(0, ci)
    cell.text = h_text
    cell.fill.solid()
    cell.fill.fore_color.rgb = PURPLE
    for p in cell.text_frame.paragraphs:
        p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = WHITE; p.font.name = FONT

for ri, (name, rule, action, purpose) in enumerate(flags):
    row = ri + 1
    for ci, val in enumerate([name, rule, action, purpose]):
        cell = tbl.cell(row, ci)
        cell.text = val
        bg = BG_LIGHT if ri % 2 == 0 else WHITE
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(14)
            p.font.color.rgb = DEEP_PURPLE if ci == 0 else TEXT_BODY
            p.font.bold = ci == 0
            p.font.name = FONT

set_footer(slide)
add_gt_to_slide(slide)


# ══════════════════════════════════════════════════════════════
# SLIDE 8: SIMULATION PIPELINE
# ══════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[2])
for ph in slide.placeholders:
    idx = ph.placeholder_format.idx
    if idx == 11:
        ph.text = "ARCHITECTURE"
        for p in ph.text_frame.paragraphs:
            p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = PURPLE; p.font.name = FONT
    elif idx == 0:
        ph.text = "Every generated spec maps to real Aera components"
        for p in ph.text_frame.paragraphs:
            p.font.size = Pt(28); p.font.bold = True; p.font.color.rgb = BLACK; p.font.name = FONT
    elif idx == 10:
        ph.text = "The simulation pipeline produces four artifact types, each validated against 21 UI components and 22 Process Builder nodes bundled in the engine."
        for p in ph.text_frame.paragraphs:
            p.font.size = Pt(18); p.font.color.rgb = BLACK; p.font.name = FONT

# Pattern B: 2x2 grid
gcw = (CW - 0.20) / 2
gch = (AH - 0.85) / 2
artifacts = [
    ("Mermaid decision flows", ["Validated flowchart syntax", "Aera component labels", "Retry up to 3x on validation failure"], svg_icon_document()),
    ("YAML component maps", ["Maps to Streams, Cortex, PB, Agent Teams", "confirmed/inferred confidence flags", "Knowledge base enforcement (KNOW-04)"], svg_icon_layers()),
    ("Mock decision tests", ["Sample inputs from actual client financials", "Uses decision_articulation from L4 data", "Expected outputs with tolerance bands"], svg_icon_check()),
    ("Integration surfaces", ["Source systems from company ERP stack", "Aera pipeline mapping", "Unmatched sources marked as TBD"], svg_icon_network()),
]
for i, (title, bullets, icon) in enumerate(artifacts):
    col = i % 2
    row = i // 2
    cx = ML + col * (gcw + 0.20)
    cy = CY + 0.55 + row * (gch + 0.12)
    _add_card(slide, cx, cy, gcw, gch, title, bullets, icon)

set_footer(slide)
add_gt_to_slide(slide)


# ══════════════════════════════════════════════════════════════
# SLIDE 9: THE CHATFN PATTERN
# ══════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[2])
for ph in slide.placeholders:
    idx = ph.placeholder_format.idx
    if idx == 11:
        ph.text = "ARCHITECTURE"
        for p in ph.text_frame.paragraphs:
            p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = PURPLE; p.font.name = FONT
    elif idx == 0:
        ph.text = "One interface, zero scoring code changes"
        for p in ph.text_frame.paragraphs:
            p.font.size = Pt(28); p.font.bold = True; p.font.color.rgb = BLACK; p.font.name = FONT
    elif idx == 10:
        ph.text = "The ChatFn dependency injection pattern lets the same scoring logic run on local Ollama or cloud vLLM -- no code changes to scoring, triage, or simulation."
        for p in ph.text_frame.paragraphs:
            p.font.size = Pt(18); p.font.color.rgb = BLACK; p.font.name = FONT

# Pattern J: Label + content rows
LABEL_W = 1.80
SEP_H = 0.08
GAP = 0.10
num_rows = 3
row_h = (AH - 0.55 - SEP_H * (num_rows - 1) - GAP * (num_rows - 1)) / num_rows

rows_data = [
    ("ChatFn\nInterface", "type ChatFn = (messages, format) => Promise<ChatResult>", "The seam that decouples scoring from transport. Three lens scorers, the retry policy, and simulation generators all accept ChatFn via dependency injection."),
    ("Ollama\nAdapter", "ollamaChat(messages, format) --> POST /api/chat", "Default path for offline operation. Runs on 36GB Apple Silicon with Qwen 8B (triage) and 32B (scoring). Zero cloud dependency."),
    ("vLLM\nAdapter", "createVllmChatFn(baseUrl, model) --> POST /v1/chat/completions", "Cloud path translating Ollama format to OpenAI response_format. Schema translator handles $ref resolution and additionalProperties stripping."),
]

y_pos = CY + 0.55
for i, (label_text, code, desc) in enumerate(rows_data):
    # Label block
    lbl = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(ML), Inches(y_pos), Inches(LABEL_W), Inches(row_h))
    lbl.fill.solid()
    lbl.fill.fore_color.rgb = DEEP_PURPLE
    lbl.line.fill.background()
    lbl_tb = slide.shapes.add_textbox(
        Inches(ML), Inches(y_pos), Inches(LABEL_W), Inches(row_h))
    lbl_tf = lbl_tb.text_frame
    lbl_tf.word_wrap = True
    lbl_tf.margin_top = Inches(row_h * 0.30)
    p = lbl_tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = label_text; r.font.size = Pt(14); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = FONT

    # Content card
    card_x = ML + LABEL_W + 0.06
    card_w = CW - LABEL_W - 0.06
    card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(card_x), Inches(y_pos), Inches(card_w), Inches(row_h))
    card.fill.solid()
    card.fill.fore_color.rgb = BG_LIGHT
    card.line.color.rgb = LIGHT_GRAY
    card.line.width = Pt(0.75)

    tb = slide.shapes.add_textbox(
        Inches(card_x + 0.15), Inches(y_pos + 0.10),
        Inches(card_w - 0.30), Inches(row_h - 0.20))
    tf = tb.text_frame; tf.word_wrap = True
    p0 = tf.paragraphs[0]
    r = p0.add_run(); r.text = code; r.font.size = Pt(14); r.font.bold = True; r.font.color.rgb = DEEP_PURPLE; r.font.name = FONT
    p1 = tf.add_paragraph(); p1.space_before = Pt(6)
    r = p1.add_run(); r.text = desc; r.font.size = Pt(14); r.font.color.rgb = TEXT_BODY; r.font.name = FONT

    y_pos += row_h + GAP
    if i < num_rows - 1:
        add_svg_native(slide, svg_divider_gradient(w=800), x=ML, y=y_pos - GAP/2, w=CW, h=0.06)

set_footer(slide)
add_gt_to_slide(slide)


# ══════════════════════════════════════════════════════════════
# SLIDE 10: THE AGENTIC DEVELOPMENT MODEL
# ══════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[2])
for ph in slide.placeholders:
    idx = ph.placeholder_format.idx
    if idx == 11:
        ph.text = "HOW WE BUILT IT"
        for p in ph.text_frame.paragraphs:
            p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = PURPLE; p.font.name = FONT
    elif idx == 0:
        ph.text = "Claude Code as autonomous builder, not a copilot"
        for p in ph.text_frame.paragraphs:
            p.font.size = Pt(28); p.font.bold = True; p.font.color.rgb = BLACK; p.font.name = FONT
    elif idx == 10:
        ph.text = "Not suggesting snippets. Executing 31 atomic plans with TDD, git commits, and checkpoint recovery -- fully autonomous within each plan boundary."
        for p in ph.text_frame.paragraphs:
            p.font.size = Pt(18); p.font.color.rgb = BLACK; p.font.name = FONT

# Pattern E: Three-column comparison
tcw = (CW - 2 * 0.15) / 3
tc_h = AH - 0.55

cols = [
    ("Traditional copilot", svg_icon_code(),
     ["Developer writes code", "AI suggests completions", "Human reviews each snippet", "Linear, manual process", "Speed gain: 10-30%"]),
    ("Agentic development", svg_icon_brain(),
     ["Human defines plan (PLAN.md)", "Agent writes tests first (TDD)", "Agent implements until green", "Agent commits and verifies", "Speed gain: 10-50x per plan"]),
    ("What changed", svg_icon_rocket(),
     ["Human role: strategy + planning", "Agent role: implementation", "Plans are the unit of work", "STATE.md enables continuity", "93 min for 213K LOC"]),
]

for i, (title, icon, bullets) in enumerate(cols):
    cx = ML + i * (tcw + 0.15)
    cy = CY + 0.55
    _add_card(slide, cx, cy, tcw, tc_h, title, bullets, icon)

set_footer(slide)
add_gt_to_slide(slide)


# ══════════════════════════════════════════════════════════════
# SLIDE 11: FINE-GRAINED PLANNING
# ══════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[2])
for ph in slide.placeholders:
    idx = ph.placeholder_format.idx
    if idx == 11:
        ph.text = "HOW WE BUILT IT"
        for p in ph.text_frame.paragraphs:
            p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = PURPLE; p.font.name = FONT
    elif idx == 0:
        ph.text = "31 atomic plans across 11 phases, 3 min average"
        for p in ph.text_frame.paragraphs:
            p.font.size = Pt(28); p.font.bold = True; p.font.color.rgb = BLACK; p.font.name = FONT
    elif idx == 10:
        ph.text = "Each plan has a clear spec (PLAN.md), verified output (SUMMARY.md + VERIFICATION.md), and an atomic git commit. Plans are the unit of resumable work."
        for p in ph.text_frame.paragraphs:
            p.font.size = Pt(18); p.font.color.rgb = BLACK; p.font.name = FONT

# Layers diagram showing the planning hierarchy
layers_svg = svg_layers_diagram(
    ["PROJECT.md -- charter + constraints",
     "ROADMAP.md -- 11 phases + requirements",
     "PLAN.md -- per-task spec (3-5 min scope)",
     "SUMMARY.md -- verified output + decisions",
     "STATE.md -- checkpoint for next session"],
    w=700, h=280)
add_svg_native(slide, layers_svg, x=ML + 0.50, y=CY + 0.55, w=8.0, h=3.80)

# Side annotation
_add_body_text(slide, ML + 9.0, CY + 0.70, 3.50, 4.00,
    ["Each plan boundary:",
     "  1. Read STATE.md",
     "  2. Execute PLAN.md with TDD",
     "  3. Write SUMMARY.md",
     "  4. Git commit",
     "  5. Update STATE.md",
     "",
     "3-minute plans are the sweet",
     "spot: context stays fresh and",
     "errors stay small."],
    font_size=Pt(14), color=TEXT_BODY)

set_footer(slide)
add_gt_to_slide(slide)


# ══════════════════════════════════════════════════════════════
# SLIDE 12: THE SESSION LOOP (run.sh)
# ══════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[2])
for ph in slide.placeholders:
    idx = ph.placeholder_format.idx
    if idx == 11:
        ph.text = "HOW WE BUILT IT"
        for p in ph.text_frame.paragraphs:
            p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = PURPLE; p.font.name = FONT
    elif idx == 0:
        ph.text = "run.sh overcomes context window limits"
        for p in ph.text_frame.paragraphs:
            p.font.size = Pt(28); p.font.bold = True; p.font.color.rgb = BLACK; p.font.name = FONT
    elif idx == 10:
        ph.text = "An infinite bash loop restarts Claude Code on context limits. Each new session reads STATE.md and picks up where the last session stopped."
        for p in ph.text_frame.paragraphs:
            p.font.size = Pt(18); p.font.color.rgb = BLACK; p.font.name = FONT

# Circular flow diagram using chevrons
flow_steps = ["Launch\nClaude Code", "Execute\nPlan", "Git\nCommit", "Context\nLimit", "Restart +\nRead STATE"]
flow_svg = svg_chevron_flow(["Launch Claude", "Execute Plan", "Git Commit", "Context Limit", "Read STATE"], w=800, h=60)
add_svg_native(slide, flow_svg, x=ML, y=CY + 0.60, w=CW, h=0.55)

# Add refresh icon to show the loop
add_svg_native(slide, svg_icon_refresh(), x=ML + CW/2 - 0.25, y=CY + 1.30, w=0.50, h=0.50)

# Key details in accent-bar cards
card_data = [
    ("Session 1", "Reads program.md and STATE.md. Executes plans until context limit. Commits progress."),
    ("Session 2+", "Reads next-session.md handoff note. Picks up from last completed plan. No work is repeated."),
    ("Resilience", "10 sessions over 2 days shipped the entire v1.0 milestone. Each session is self-contained."),
]
card_y = CY + 1.95
card_h = (BY - card_y - 0.10) / 3
for i, (title, desc) in enumerate(card_data):
    cy = card_y + i * (card_h + 0.08)
    # Accent bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(ML), Inches(cy), Inches(0.06), Inches(card_h - 0.04))
    bar.fill.solid(); bar.fill.fore_color.rgb = PURPLE; bar.line.fill.background()
    # Card bg
    cbg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(ML + 0.06), Inches(cy), Inches(CW - 0.06), Inches(card_h - 0.04))
    cbg.fill.solid(); cbg.fill.fore_color.rgb = WHITE; cbg.line.color.rgb = LIGHT_GRAY; cbg.line.width = Pt(0.75)
    # Text
    ttb = slide.shapes.add_textbox(Inches(ML + 0.25), Inches(cy + 0.08), Inches(2.00), Inches(0.35))
    tf = ttb.text_frame; tf.word_wrap = True
    r = tf.paragraphs[0].add_run(); r.text = title; r.font.size = Pt(14); r.font.bold = True; r.font.color.rgb = DEEP_PURPLE; r.font.name = FONT
    dtb = slide.shapes.add_textbox(Inches(ML + 2.40), Inches(cy + 0.08), Inches(CW - 2.65), Inches(card_h - 0.20))
    df = dtb.text_frame; df.word_wrap = True
    r = df.paragraphs[0].add_run(); r.text = desc; r.font.size = Pt(14); r.font.color.rgb = TEXT_BODY; r.font.name = FONT

set_footer(slide)
add_gt_to_slide(slide)


# ══════════════════════════════════════════════════════════════
# SLIDE 13: TDD AT SCALE
# ══════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[2])
for ph in slide.placeholders:
    idx = ph.placeholder_format.idx
    if idx == 11:
        ph.text = "HOW WE BUILT IT"
        for p in ph.text_frame.paragraphs:
            p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = PURPLE; p.font.name = FONT
    elif idx == 0:
        ph.text = "412 tests, zero external dependencies required"
        for p in ph.text_frame.paragraphs:
            p.font.size = Pt(28); p.font.bold = True; p.font.color.rgb = BLACK; p.font.name = FONT
    elif idx == 10:
        ph.text = "Node.js built-in test runner with pure-function core logic. Tests run without Ollama -- deterministic by design via dependency injection."
        for p in ph.text_frame.paragraphs:
            p.font.size = Pt(18); p.font.color.rgb = BLACK; p.font.name = FONT

# Metric table showing test progression
tbl = slide.shapes.add_table(7, 3, Inches(ML), Inches(CY + 0.55), Inches(CW), Inches(4.30)).table
tbl.columns[0].width = Inches(5.00)
tbl.columns[1].width = Inches(3.75)
tbl.columns[2].width = Inches(3.75)

h_data = ["Phase", "Tests Added", "Pattern"]
for ci, h in enumerate(h_data):
    cell = tbl.cell(0, ci)
    cell.text = h
    cell.fill.solid(); cell.fill.fore_color.rgb = PURPLE
    for p in cell.text_frame.paragraphs:
        p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = WHITE; p.font.name = FONT

rows_data = [
    ("Phase 1: Foundation", "22", "Zod schema validation"),
    ("Phase 3: Triage + Red Flags", "61", "Pure-function triage with 5 red flags"),
    ("Phase 4: Scoring Engine", "48", "DI via chatFn -- tested without Ollama"),
    ("Phase 6: Simulation", "62", "Mermaid/YAML validation + knowledge base"),
    ("Phase 8: Resilience", "56", "Checkpoint recovery + retry policies"),
    ("Phase 9: Reports", "42", "Markdown/TSV output formatting"),
]
for ri, (phase, count, pattern) in enumerate(rows_data):
    bg = BG_LIGHT if ri % 2 == 0 else WHITE
    for ci, val in enumerate([phase, count, pattern]):
        cell = tbl.cell(ri + 1, ci)
        cell.text = val
        cell.fill.solid(); cell.fill.fore_color.rgb = bg
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(14)
            p.font.color.rgb = DEEP_PURPLE if ci == 0 else TEXT_BODY
            p.font.bold = ci == 0
            p.font.name = FONT

set_footer(slide)
add_gt_to_slide(slide)


# ══════════════════════════════════════════════════════════════
# SLIDE 14: NO FRAMEWORKS
# ══════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[2])
for ph in slide.placeholders:
    idx = ph.placeholder_format.idx
    if idx == 11:
        ph.text = "HOW WE BUILT IT"
        for p in ph.text_frame.paragraphs:
            p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = PURPLE; p.font.name = FONT
    elif idx == 0:
        ph.text = "150 LOC of async generators replace 10,000+ LOC"
        for p in ph.text_frame.paragraphs:
            p.font.size = Pt(28); p.font.bold = True; p.font.color.rgb = BLACK; p.font.name = FONT
    elif idx == 10:
        ph.text = "The orchestrator.ts is a hand-written agentic loop using standard async generators. No LangChain, CrewAI, or Autogen -- proving robust patterns work with stdlib."
        for p in ph.text_frame.paragraphs:
            p.font.size = Pt(18); p.font.color.rgb = BLACK; p.font.name = FONT

# Pattern D: Two-column comparison
col_w = (CW - 0.20) / 2
col_h = AH - 0.60

_add_card(slide, ML, CY + 0.55, col_w, col_h, "Our orchestrator (~150 LOC)",
    ["async function* runAgent(history)",
     "Yields SSE events via AsyncGenerator",
     "Tool calls: getTool() from registry",
     "Streaming: provider.stream(history)",
     "Max rounds limit with forced summary",
     "Result truncation prevents context explosion",
     "",
     "Dependencies: 0 framework libraries"],
    svg_icon_code())

_add_card(slide, ML + col_w + 0.20, CY + 0.55, col_w, col_h, "Framework alternatives",
    ["LangChain: ~10,000+ LOC in agent modules",
     "CrewAI: multi-agent abstractions",
     "Autogen: conversation orchestration",
     "",
     "Common trade-offs:",
     "  Heavy dependency trees",
     "  Opinionated abstractions",
     "  Version compatibility issues",
     "  Harder to debug and customize"],
    svg_icon_gear())

set_footer(slide)
add_gt_to_slide(slide)


# ══════════════════════════════════════════════════════════════
# SLIDE 15: V1.0 METRICS DASHBOARD
# ══════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[2])
for ph in slide.placeholders:
    idx = ph.placeholder_format.idx
    if idx == 11:
        ph.text = "RESULTS"
        for p in ph.text_frame.paragraphs:
            p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = PURPLE; p.font.name = FONT
    elif idx == 0:
        ph.text = "v1.0 shipped in 2 days with 100% requirement coverage"
        for p in ph.text_frame.paragraphs:
            p.font.size = Pt(28); p.font.bold = True; p.font.color.rgb = BLACK; p.font.name = FONT
    elif idx == 10:
        ph.text = "These are evidence metrics from real project artifacts -- every number is traceable to git commits, test outputs, and planning documents."
        for p in ph.text_frame.paragraphs:
            p.font.size = Pt(18); p.font.color.rgb = BLACK; p.font.name = FONT

metrics_svg = svg_metrics_dashboard()
add_svg_native(slide, metrics_svg, x=ML, y=CY + 0.55, w=CW, h=3.80)

# Source citation
src_tb = slide.shapes.add_textbox(Inches(ML), Inches(BY - 0.35), Inches(CW), Inches(0.30))
sf = src_tb.text_frame; sf.word_wrap = True
p = sf.paragraphs[0]
r = p.add_run(); r.text = "Source: .planning/MILESTONES.md, .planning/RETROSPECTIVE.md -- all metrics from actual project artifacts"
r.font.size = Pt(14); r.font.italic = True; r.font.color.rgb = TEXT_SUB; r.font.name = FONT

set_footer(slide)
add_gt_to_slide(slide)


# ══════════════════════════════════════════════════════════════
# SLIDE 16: V1.1 CLOUD ACCELERATION
# ══════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[2])
for ph in slide.placeholders:
    idx = ph.placeholder_format.idx
    if idx == 11:
        ph.text = "RESULTS"
        for p in ph.text_frame.paragraphs:
            p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = PURPLE; p.font.name = FONT
    elif idx == 0:
        ph.text = "17 hours local to under 30 minutes cloud"
        for p in ph.text_frame.paragraphs:
            p.font.size = Pt(28); p.font.bold = True; p.font.color.rgb = BLACK; p.font.name = FONT
    elif idx == 10:
        ph.text = "v1.1 adds vLLM on RunPod H100 as an optional backend. The ChatFn pattern means zero changes to scoring logic -- only the transport layer changes."
        for p in ph.text_frame.paragraphs:
            p.font.size = Pt(18); p.font.color.rgb = BLACK; p.font.name = FONT

perf_svg = svg_perf_comparison()
add_svg_native(slide, perf_svg, x=ML, y=CY + 0.55, w=CW, h=3.40)

# Key message bar
kmb_y = BY - 0.70
bar_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
    Inches(ML), Inches(kmb_y), Inches(CW), Inches(0.65))
bar_bg.fill.solid(); bar_bg.fill.fore_color.rgb = LP_BG
bar_bg.line.color.rgb = PURPLE; bar_bg.line.width = Pt(1)
kmtb = slide.shapes.add_textbox(Inches(ML + 0.20), Inches(kmb_y + 0.10), Inches(CW - 0.40), Inches(0.45))
kf = kmtb.text_frame; kf.word_wrap = True
p = kf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = "10-20 concurrent opportunities via semaphore-bounded pipeline. Pre-flight schema validation prevents wasted GPU time."
r.font.size = Pt(14); r.font.bold = True; r.font.color.rgb = DEEP_PURPLE; r.font.name = FONT

set_footer(slide)
add_gt_to_slide(slide)


# ══════════════════════════════════════════════════════════════
# SLIDE 17: OVERNIGHT RESILIENCE
# ══════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[2])
for ph in slide.placeholders:
    idx = ph.placeholder_format.idx
    if idx == 11:
        ph.text = "RESULTS"
        for p in ph.text_frame.paragraphs:
            p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = PURPLE; p.font.name = FONT
    elif idx == 0:
        ph.text = "Three tiers of resilience for unattended runs"
        for p in ph.text_frame.paragraphs:
            p.font.size = Pt(28); p.font.bold = True; p.font.color.rgb = BLACK; p.font.name = FONT
    elif idx == 10:
        ph.text = "Crash recovery, retry policies, and git auto-commit enable multi-hour scoring runs that complete without human intervention."
        for p in ph.text_frame.paragraphs:
            p.font.size = Pt(18); p.font.color.rgb = BLACK; p.font.name = FONT

# Pattern E: Three columns
tcw = (CW - 2 * 0.15) / 3
tc_h = AH - 0.55

cols = [
    ("Tier 1: Retry", svg_icon_refresh(),
     ["LLM call fails? Retry up to 3x", "Exponential backoff", "Same prompt, same parameters", "Handles transient failures"]),
    ("Tier 2: Fallback", svg_icon_gear(),
     ["Structured output fails?", "Simplify the prompt", "Reduce schema constraints", "Get partial result vs none"]),
    ("Tier 3: Skip + Log", svg_icon_shield(),
     ["Opportunity is intractable?", "Log error to checkpoint.json", "Skip and continue pipeline", "Report failures in summary"]),
]
for i, (title, icon, bullets) in enumerate(cols):
    cx = ML + i * (tcw + 0.15)
    cy = CY + 0.55
    _add_card(slide, cx, cy, tcw, tc_h, title, bullets, icon)

set_footer(slide)
add_gt_to_slide(slide)


# ══════════════════════════════════════════════════════════════
# SLIDE 18: IMPLICATIONS FOR CONSULTING
# ══════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[2])
for ph in slide.placeholders:
    idx = ph.placeholder_format.idx
    if idx == 11:
        ph.text = "IMPLICATIONS"
        for p in ph.text_frame.paragraphs:
            p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = PURPLE; p.font.name = FONT
    elif idx == 0:
        ph.text = "Agentic development is not future-state"
        for p in ph.text_frame.paragraphs:
            p.font.size = Pt(28); p.font.bold = True; p.font.color.rgb = BLACK; p.font.name = FONT
    elif idx == 10:
        ph.text = "A solo developer with Claude Code shipped production-grade software faster than traditional teams. This has direct implications for delivery timelines and staffing models."
        for p in ph.text_frame.paragraphs:
            p.font.size = Pt(18); p.font.color.rgb = BLACK; p.font.name = FONT

# Pattern O: Bullet memo
tb = slide.shapes.add_textbox(Inches(ML), Inches(CY + 0.55), Inches(CW), Inches(AH - 0.55))
tf = tb.text_frame; tf.word_wrap = True

bullets_data = [
    ("Delivery speed: ", "213K LOC in 93 minutes of automated execution challenges assumptions about project timelines. The bottleneck shifts from implementation to planning and verification."),
    ("Quality assurance: ", "412 tests with 100% requirement coverage demonstrates that agent-built code meets enterprise quality standards when guided by TDD discipline and fine-grained plans."),
    ("Staffing models: ", "One developer + one coding agent replaced what would traditionally require a team of 4-6 over 2-3 months. The human role shifts to architecture, planning, and client strategy."),
    ("Cost structure: ", "Coding agent execution costs (API usage) are a fraction of equivalent developer-hours. The economic model for software delivery changes fundamentally."),
]

for i, (bold_part, rest) in enumerate(bullets_data):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_before = Pt(12) if i > 0 else Pt(0)
    r1 = p.add_run()
    r1.text = f"  {bold_part}"
    r1.font.size = Pt(14); r1.font.bold = True; r1.font.color.rgb = PURPLE; r1.font.name = FONT
    r2 = p.add_run()
    r2.text = rest
    r2.font.size = Pt(14); r2.font.color.rgb = TEXT_BODY; r2.font.name = FONT

set_footer(slide)
add_gt_to_slide(slide)


# ══════════════════════════════════════════════════════════════
# SLIDE 19: REPLICATING THIS APPROACH
# ══════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[2])
for ph in slide.placeholders:
    idx = ph.placeholder_format.idx
    if idx == 11:
        ph.text = "IMPLICATIONS"
        for p in ph.text_frame.paragraphs:
            p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = PURPLE; p.font.name = FONT
    elif idx == 0:
        ph.text = "Four transferable patterns for any project"
        for p in ph.text_frame.paragraphs:
            p.font.size = Pt(28); p.font.bold = True; p.font.color.rgb = BLACK; p.font.name = FONT
    elif idx == 10:
        ph.text = "These patterns are not specific to this project. They work for any enterprise software engagement where a coding agent is part of the delivery team."
        for p in ph.text_frame.paragraphs:
            p.font.size = Pt(18); p.font.color.rgb = BLACK; p.font.name = FONT

# Pattern B: 2x2 grid
gcw = (CW - 0.20) / 2
gch = (AH - 0.85) / 2
patterns = [
    ("Fine-grained planning", ["3-5 minute execution windows", "PLAN.md with clear specs and exit criteria", "SUMMARY.md with verified output", "Plans are resumable units of work"], svg_icon_document()),
    ("TDD discipline", ["Tests first, implementation second", "Pure functions in core logic (no I/O)", "Dependency injection for testability", "Factory helpers for minimal test objects"], svg_icon_check()),
    ("Checkpoint + resume", ["STATE.md tracks current position", "Crash-safe: resume from last checkpoint", "Git auto-commit after each phase", "Context archive between sessions"], svg_icon_shield()),
    ("Offline-first architecture", ["Local models as default (zero cloud cost)", "Cloud as opt-in acceleration", "Same interface for both paths (ChatFn)", "No vendor lock-in on scoring logic"], svg_icon_cloud()),
]
for i, (title, bullets, icon) in enumerate(patterns):
    col = i % 2
    row = i // 2
    cx = ML + col * (gcw + 0.20)
    cy = CY + 0.55 + row * (gch + 0.12)
    _add_card(slide, cx, cy, gcw, gch, title, bullets, icon)

set_footer(slide)
add_gt_to_slide(slide)


# ══════════════════════════════════════════════════════════════
# SLIDE 20: THE SKILLS ECOSYSTEM
# ══════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[2])
for ph in slide.placeholders:
    idx = ph.placeholder_format.idx
    if idx == 11:
        ph.text = "IMPLICATIONS"
        for p in ph.text_frame.paragraphs:
            p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = PURPLE; p.font.name = FONT
    elif idx == 0:
        ph.text = "Custom skills create a composable agent toolkit"
        for p in ph.text_frame.paragraphs:
            p.font.size = Pt(28); p.font.bold = True; p.font.color.rgb = BLACK; p.font.name = FONT
    elif idx == 10:
        ph.text = "This very deck was generated by /acnpptx. Each skill is a reusable capability that compounds across projects -- the more you build, the faster you move."
        for p in ph.text_frame.paragraphs:
            p.font.size = Pt(18); p.font.color.rgb = BLACK; p.font.name = FONT

# Pattern C: Table of skills
tbl = slide.shapes.add_table(5, 3, Inches(ML), Inches(CY + 0.55), Inches(CW), Inches(3.80)).table
tbl.columns[0].width = Inches(2.50)
tbl.columns[1].width = Inches(4.50)
tbl.columns[2].width = Inches(5.50)

h_data = ["Skill", "Capability", "Impact"]
for ci, h in enumerate(h_data):
    cell = tbl.cell(0, ci)
    cell.text = h
    cell.fill.solid(); cell.fill.fore_color.rgb = PURPLE
    for p in cell.text_frame.paragraphs:
        p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = WHITE; p.font.name = FONT

skills = [
    ("/acnpptx", "PowerPoint generation with brand compliance", "This presentation -- 20 slides from project artifacts"),
    ("/gsd", "Fine-grained planning and state tracking", "31 plans, 11 phases, checkpoint recovery"),
    ("/build-e2e-skill", "End-to-end skill creation workflow", "New skills from spec to verified output"),
    ("/review-pr", "Automated pull request review", "Code quality checks on agent-generated code"),
]
for ri, (name, cap, impact) in enumerate(skills):
    bg = BG_LIGHT if ri % 2 == 0 else WHITE
    for ci, val in enumerate([name, cap, impact]):
        cell = tbl.cell(ri + 1, ci)
        cell.text = val
        cell.fill.solid(); cell.fill.fore_color.rgb = bg
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(14)
            p.font.color.rgb = DEEP_PURPLE if ci == 0 else TEXT_BODY
            p.font.bold = ci == 0
            p.font.name = FONT

# Key message bar
kmb_y = BY - 0.70
bar_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
    Inches(ML), Inches(kmb_y), Inches(CW), Inches(0.65))
bar_bg.fill.solid(); bar_bg.fill.fore_color.rgb = LP_BG
bar_bg.line.color.rgb = PURPLE; bar_bg.line.width = Pt(1)
kmtb = slide.shapes.add_textbox(Inches(ML + 0.20), Inches(kmb_y + 0.10), Inches(CW - 0.40), Inches(0.45))
kf = kmtb.text_frame; kf.word_wrap = True
p = kf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = "Skills compound: each new capability accelerates every future project that uses it"
r.font.size = Pt(14); r.font.bold = True; r.font.color.rgb = DEEP_PURPLE; r.font.name = FONT

set_footer(slide)
add_gt_to_slide(slide)


# ══════════════════════════════════════════════════════════════
# SLIDE 21: CLOSING
# ══════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[2])
for ph in slide.placeholders:
    idx = ph.placeholder_format.idx
    if idx == 11:
        ph.text = "CLOSING"
        for p in ph.text_frame.paragraphs:
            p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = PURPLE; p.font.name = FONT
    elif idx == 0:
        ph.text = "The agent-augmented consultant"
        for p in ph.text_frame.paragraphs:
            p.font.size = Pt(28); p.font.bold = True; p.font.color.rgb = BLACK; p.font.name = FONT
    elif idx == 10:
        ph.text = "The future of consulting is not AI replacing consultants. It is consultants armed with autonomous agents that handle implementation while humans focus on strategy."
        for p in ph.text_frame.paragraphs:
            p.font.size = Pt(18); p.font.color.rgb = BLACK; p.font.name = FONT

# Pattern F: Numbered summary with CTA
items = [
    ("Agents handle implementation", "Fine-grained plans + TDD + checkpoint recovery enable autonomous code delivery at enterprise scale"),
    ("Humans focus on strategy", "Client relationships, adoption planning, architecture decisions, and quality verification remain human responsibilities"),
    ("The toolkit is real today", "Claude Code + custom skills + the patterns shown here are available now -- not a future roadmap item"),
]

item_h = 1.20
item_y = CY + 0.55
for i, (title, desc) in enumerate(items):
    iy = item_y + i * (item_h + 0.12)

    # Number badge
    add_svg_native(slide, svg_callout_badge(str(i + 1), size=40), x=ML, y=iy + 0.15, w=0.40, h=0.40)

    # Card
    card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(ML + 0.55), Inches(iy), Inches(CW - 0.55), Inches(item_h))
    card.fill.solid(); card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = LIGHT_GRAY; card.line.width = Pt(0.75)

    ttb = slide.shapes.add_textbox(Inches(ML + 0.70), Inches(iy + 0.10), Inches(CW - 0.85), Inches(0.35))
    tf = ttb.text_frame; tf.word_wrap = True
    r = tf.paragraphs[0].add_run(); r.text = title; r.font.size = Pt(14); r.font.bold = True; r.font.color.rgb = DEEP_PURPLE; r.font.name = FONT

    dtb = slide.shapes.add_textbox(Inches(ML + 0.70), Inches(iy + 0.45), Inches(CW - 0.85), Inches(item_h - 0.55))
    df = dtb.text_frame; df.word_wrap = True
    r = df.paragraphs[0].add_run(); r.text = desc; r.font.size = Pt(14); r.font.color.rgb = TEXT_BODY; r.font.name = FONT

# CTA bar
cta_y = item_y + 3 * (item_h + 0.12) + 0.10
cta = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
    Inches(ML), Inches(cta_y), Inches(CW), Inches(0.70))
cta.fill.solid(); cta.fill.fore_color.rgb = GOLD
cta.line.fill.background()
cta_tb = slide.shapes.add_textbox(Inches(ML + 0.20), Inches(cta_y + 0.10), Inches(CW - 0.40), Inches(0.50))
cf = cta_tb.text_frame; cf.word_wrap = True
p = cf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = "Ready to explore agentic development for your next engagement? Let's talk."
r.font.size = Pt(16); r.font.bold = True; r.font.color.rgb = BLACK; r.font.name = FONT

set_footer(slide)
add_gt_to_slide(slide)


# ══════════════════════════════════════════════════════════════
# SAVE + VERIFY
# ══════════════════════════════════════════════════════════════

output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "agent-factory-consultant-deck.pptx")
prs.save(output_path)
cleanup_temp()
print(f"Saved: {output_path}")
print(f"Slides: {len(prs.slides)}")

import subprocess
_venv_py = os.path.join(os.path.dirname(SKILL_DIR), ".venv", "Scripts", "python.exe")
_py = _venv_py if os.path.exists(_venv_py) else sys.executable

# Structural check
print("\n=== Structural Verification ===")
subprocess.run([_py, os.path.join(SKILL_DIR, "scripts", "verify_pptx.py"), output_path], check=False)

# Thumbnail export
thumb_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)), "thumbnails")
subprocess.run([_py, os.path.join(SKILL_DIR, "scripts", "thumbnail.py"), output_path, thumb_dir], check=False)

print(f"\nDone. Output: {output_path}")
