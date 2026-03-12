"""
Aera Skill Feasibility Engine — Capability Briefing Deck
18 slides: methodology briefing for management consultants and technology advisors.

Patterns used: Cover(0), D, E, C, M, J, R, I, C, E, B, V, E, F, D, G, O, Cover(6)
"""

import sys
import os

# ── Skill path setup ──────────────────────────────────────────
def _find_skill_scripts():
    skills_dir = os.path.join(os.path.expanduser("~"), ".claude", "skills")
    for candidate in ["acnpptx", "pptx"]:
        path = os.path.join(skills_dir, candidate, "scripts")
        if os.path.exists(os.path.join(path, "helpers.py")):
            return path
    raise FileNotFoundError("Cannot find acnpptx skill scripts.")

sys.path.insert(0, _find_skill_scripts())
from helpers import *
from svg_pipeline import add_svg_native as _add_svg_native, cleanup_temp
from pptx import Presentation
from pptx.oxml.ns import qn as _qn
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ── Output path ───────────────────────────────────────────────
OUTPUT = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                      "feasibility-engine-capability-deck.pptx")

# ── Load template and clear slides ────────────────────────────
prs = Presentation(TEMPLATE_PATH)
while len(prs.slides) > 0:
    sldId = prs.slides._sldIdLst[0]
    rId = sldId.get(_qn('r:id'))
    prs.part.drop_rel(rId)
    del prs.slides._sldIdLst[0]

reset_slide_counter()

def add_svg_native(slide, svg_str, x, y, w, h, png_width=800):
    return _add_svg_native(slide, prs, svg_str, x, y, w, h, png_width)


# ══════════════════════════════════════════════════════════════
# SVG helper functions
# ══════════════════════════════════════════════════════════════

def svg_grouped_bar_chart(categories, series1, series2, label1, label2,
                          w=700, h=260, c1="#D8D8D8", c2="#7E00FF"):
    """Grouped bar chart comparing two series."""
    pad_l, pad_r, pad_t, pad_b = 55, 20, 30, 50
    plot_w = w - pad_l - pad_r
    plot_h = h - pad_t - pad_b
    n = len(categories)
    group_w = plot_w / n
    bar_w = group_w * 0.32
    gap = group_w * 0.06
    max_val = max(max(series1), max(series2))
    bars = ""
    for i in range(n):
        gx = pad_l + i * group_w
        # Series 1
        h1 = (series1[i] / max_val) * plot_h
        y1 = pad_t + plot_h - h1
        bx1 = gx + gap
        bars += f'<rect x="{bx1}" y="{y1}" width="{bar_w}" height="{h1}" fill="{c1}" rx="2"/>'
        bars += (f'<text x="{bx1 + bar_w/2}" y="{y1 - 6}" text-anchor="middle" '
                 f'font-family="Graphik, sans-serif" font-size="13" fill="#333333" '
                 f'font-weight="bold">{series1[i]}%</text>')
        # Series 2
        h2 = (series2[i] / max_val) * plot_h
        y2 = pad_t + plot_h - h2
        bx2 = gx + gap + bar_w + gap
        bars += f'<rect x="{bx2}" y="{y2}" width="{bar_w}" height="{h2}" fill="{c2}" rx="2"/>'
        bars += (f'<text x="{bx2 + bar_w/2}" y="{y2 - 6}" text-anchor="middle" '
                 f'font-family="Graphik, sans-serif" font-size="13" fill="#333333" '
                 f'font-weight="bold">{series2[i]}%</text>')
        # Category label
        bars += (f'<text x="{gx + group_w/2}" y="{h - 8}" text-anchor="middle" '
                 f'font-family="Graphik, sans-serif" font-size="12" fill="#666666">{categories[i]}</text>')
    # Axis
    bars += (f'<line x1="{pad_l}" y1="{pad_t + plot_h}" x2="{pad_l + plot_w}" '
             f'y2="{pad_t + plot_h}" stroke="#D8D8D8" stroke-width="1"/>')
    # Legend
    lx = pad_l + 10
    ly = pad_t - 14
    bars += f'<rect x="{lx}" y="{ly}" width="12" height="12" fill="{c1}" rx="2"/>'
    bars += (f'<text x="{lx + 16}" y="{ly + 10}" font-family="Graphik, sans-serif" '
             f'font-size="12" fill="#666666">{label1}</text>')
    lx2 = lx + 130
    bars += f'<rect x="{lx2}" y="{ly}" width="12" height="12" fill="{c2}" rx="2"/>'
    bars += (f'<text x="{lx2 + 16}" y="{ly + 10}" font-family="Graphik, sans-serif" '
             f'font-size="12" fill="#666666">{label2}</text>')
    return (f'<svg xmlns="http://www.w3.org/2000/svg" width="{w}" height="{h}">'
            f'{bars}</svg>')


def svg_stacked_horiz_bar(segments, labels, total_w=700, h=55,
                          colors=None):
    """Single stacked horizontal bar with labeled segments."""
    if colors is None:
        colors = ["#5F0095", "#7E00FF", "#D8D8D8"]
    pad_l = 0
    bar_h = 36
    bar_y = 4
    total = sum(s for s in segments)
    bars = ""
    x = pad_l
    for i, (seg, lbl) in enumerate(zip(segments, labels)):
        sw = (seg / total) * total_w
        c = colors[i % len(colors)]
        bars += f'<rect x="{x}" y="{bar_y}" width="{sw}" height="{bar_h}" fill="{c}" rx="3"/>'
        # Label inside
        tc = "#FFFFFF" if c != "#D8D8D8" else "#333333"
        bars += (f'<text x="{x + sw/2}" y="{bar_y + bar_h/2 + 1}" text-anchor="middle" '
                 f'dominant-baseline="middle" font-family="Graphik, sans-serif" font-size="13" '
                 f'fill="{tc}" font-weight="bold">{lbl}</text>')
        x += sw
    return (f'<svg xmlns="http://www.w3.org/2000/svg" width="{total_w}" height="{h}">'
            f'{bars}</svg>')


def svg_action_matrix(w=700, h=200):
    """Red flag action matrix table as SVG."""
    rows = [
        ("DEAD_ZONE", "No decisions to automate", "SKIP"),
        ("PHANTOM", "Not recognized by the business", "SKIP"),
        ("NO_STAKES", "No financial impact", "DEMOTE"),
        ("CONFIDENCE_GAP", "Low organizational conviction", "FLAG"),
        ("ORPHAN", "Incomplete definition", "FLAG"),
    ]
    row_h = h / (len(rows) + 1)
    hdr_h = row_h
    col_ws = [0.22 * w, 0.52 * w, 0.26 * w]
    svg = ""
    # Header
    svg += f'<rect x="0" y="0" width="{w}" height="{hdr_h}" fill="#5F0095" rx="3"/>'
    hdrs = ["Red Flag", "Detection Criterion", "Action"]
    cx = 0
    for i, (hdr, cw) in enumerate(zip(hdrs, col_ws)):
        svg += (f'<text x="{cx + cw/2}" y="{hdr_h/2 + 1}" text-anchor="middle" '
                f'dominant-baseline="middle" font-family="Graphik, sans-serif" '
                f'font-size="13" fill="white" font-weight="bold">{hdr}</text>')
        cx += cw
    # Rows
    for ri, (flag, crit, action) in enumerate(rows):
        ry = hdr_h + ri * row_h
        bg = "#F5F5F5" if ri % 2 == 0 else "#FFFFFF"
        svg += f'<rect x="0" y="{ry}" width="{w}" height="{row_h}" fill="{bg}"/>'
        # Flag name
        svg += (f'<text x="{col_ws[0]/2}" y="{ry + row_h/2 + 1}" text-anchor="middle" '
                f'dominant-baseline="middle" font-family="Graphik, sans-serif" '
                f'font-size="12" fill="#5F0095" font-weight="bold">{flag}</text>')
        # Criterion
        svg += (f'<text x="{col_ws[0] + 8}" y="{ry + row_h/2 + 1}" '
                f'dominant-baseline="middle" font-family="Graphik, sans-serif" '
                f'font-size="12" fill="#333333">{crit}</text>')
        # Action badge
        ac = "#7E00FF" if action == "SKIP" else ("#5F0095" if action == "DEMOTE" else "#D8D8D8")
        tc = "#FFFFFF" if action != "FLAG" else "#333333"
        ax = col_ws[0] + col_ws[1] + col_ws[2] / 2
        bw = 70
        svg += (f'<rect x="{ax - bw/2}" y="{ry + row_h/2 - 11}" width="{bw}" '
                f'height="22" rx="11" fill="{ac}"/>')
        svg += (f'<text x="{ax}" y="{ry + row_h/2 + 1}" text-anchor="middle" '
                f'dominant-baseline="middle" font-family="Graphik, sans-serif" '
                f'font-size="11" fill="{tc}" font-weight="bold">{action}</text>')
    return (f'<svg xmlns="http://www.w3.org/2000/svg" width="{w}" height="{h}">'
            f'{svg}</svg>')


def svg_timeline_compare(w=700, h=180):
    """Local vs cloud timeline comparison SVG."""
    pad_l, pad_t = 110, 20
    bar_h = 40
    gap = 30
    max_w = w - pad_l - 30
    # Local: 17 hours
    local_w = max_w
    local_y = pad_t
    cloud_w = max_w * (0.5 / 17.0)
    cloud_y = pad_t + bar_h + gap

    svg = ""
    # Local bar
    svg += f'<rect x="{pad_l}" y="{local_y}" width="{local_w}" height="{bar_h}" fill="#D8D8D8" rx="4"/>'
    svg += (f'<text x="{pad_l - 8}" y="{local_y + bar_h/2 + 1}" text-anchor="end" '
            f'dominant-baseline="middle" font-family="Graphik, sans-serif" '
            f'font-size="14" fill="#333333" font-weight="bold">Local (Ollama)</text>')
    svg += (f'<text x="{pad_l + local_w/2}" y="{local_y + bar_h/2 + 1}" text-anchor="middle" '
            f'dominant-baseline="middle" font-family="Graphik, sans-serif" '
            f'font-size="14" fill="#333333" font-weight="bold">~17 hours</text>')
    # Cloud bar
    svg += f'<rect x="{pad_l}" y="{cloud_y}" width="{cloud_w}" height="{bar_h}" fill="#7E00FF" rx="4"/>'
    svg += (f'<text x="{pad_l - 8}" y="{cloud_y + bar_h/2 + 1}" text-anchor="end" '
            f'dominant-baseline="middle" font-family="Graphik, sans-serif" '
            f'font-size="14" fill="#333333" font-weight="bold">Cloud (vLLM)</text>')
    svg += (f'<text x="{pad_l + cloud_w + 10}" y="{cloud_y + bar_h/2 + 1}" text-anchor="start" '
            f'dominant-baseline="middle" font-family="Graphik, sans-serif" '
            f'font-size="14" fill="#7E00FF" font-weight="bold">&lt;30 min</text>')
    # Labels
    note_y = cloud_y + bar_h + 25
    svg += (f'<text x="{pad_l}" y="{note_y}" font-family="Graphik, sans-serif" '
            f'font-size="12" fill="#666666">339-opportunity catalog  |  Same scoring logic, same artifacts  |  &lt;$10/run</text>')
    return (f'<svg xmlns="http://www.w3.org/2000/svg" width="{w}" height="{h}">'
            f'{svg}</svg>')


def svg_chevron_flow(items, w=800, h=70, color="#7E00FF"):
    n = len(items)
    cw = w / n
    arrow = 14
    shapes = ""
    for i, lbl in enumerate(items):
        x = i * cw
        op = max(0.5, 1.0 - i * 0.12)
        shapes += (f'<polygon points="{x},0 {x+cw-arrow},0 {x+cw},{h/2} '
                   f'{x+cw-arrow},{h} {x},{h}" fill="{color}" opacity="{op}"/>')
        shapes += (f'<text x="{x+cw/2}" y="{h/2+5}" text-anchor="middle" '
                   f'font-family="Graphik, sans-serif" font-size="13" fill="white" '
                   f'font-weight="bold">{lbl}</text>')
    return (f'<svg xmlns="http://www.w3.org/2000/svg" width="{w}" height="{h}">'
            f'{shapes}</svg>')


def svg_callout_badge(text="1", size=40, bg_color="#5F0095", text_color="#FFFFFF"):
    return (f'<svg xmlns="http://www.w3.org/2000/svg" width="{size}" height="{size}" '
            f'viewBox="0 0 {size} {size}">'
            f'<circle cx="{size//2}" cy="{size//2}" r="{size//2-2}" fill="{bg_color}"/>'
            f'<text x="{size//2}" y="{size//2+1}" text-anchor="middle" dominant-baseline="central" '
            f'font-family="Graphik, sans-serif" font-size="{int(size*0.45)}" font-weight="bold" '
            f'fill="{text_color}">{text}</text></svg>')


def svg_highlight_bar(w=800, h=60, bg_color="#F3E8FF", border_color="#7E00FF"):
    return (f'<svg xmlns="http://www.w3.org/2000/svg" width="{w}" height="{h}" viewBox="0 0 {w} {h}">'
            f'<rect x="0" y="0" width="{w}" height="{h}" rx="4" fill="{bg_color}" '
            f'stroke="{border_color}" stroke-width="2"/>'
            f'<polygon points="4,{h//2-8} 20,{h//2} 4,{h//2+8}" fill="{border_color}"/></svg>')


def svg_divider_gradient(w=800, color="#7E00FF"):
    return (f'<svg xmlns="http://www.w3.org/2000/svg" width="{w}" height="6" viewBox="0 0 {w} 6">'
            f'<defs><linearGradient id="dg" x1="0%" y1="0%" x2="100%" y2="0%">'
            f'<stop offset="0%" stop-color="{color}" stop-opacity="1"/>'
            f'<stop offset="100%" stop-color="{color}" stop-opacity="0.1"/>'
            f'</linearGradient></defs>'
            f'<rect x="0" y="0" width="{w}" height="6" rx="3" fill="url(#dg)"/></svg>')


def svg_layers_diagram(layers, w=500, h=240, color="#7E00FF"):
    n = len(layers)
    lh = (h - (n - 1) * 6) / n
    shapes = ""
    for i, lbl in enumerate(layers):
        y = i * (lh + 6)
        op = max(0.35, 1.0 - i * 0.18)
        shapes += (f'<rect x="20" y="{y}" width="{w-40}" height="{lh}" rx="6" '
                   f'fill="{color}" opacity="{op}"/>')
        shapes += (f'<text x="{w/2}" y="{y+lh/2+5}" text-anchor="middle" '
                   f'font-family="Graphik, sans-serif" font-size="14" fill="white" '
                   f'font-weight="bold">{lbl}</text>')
    return (f'<svg xmlns="http://www.w3.org/2000/svg" width="{w}" height="{int(n*(lh+6)-6)}">'
            f'{shapes}</svg>')


def svg_icon_shield(size=64, color="#7E00FF"):
    return f'''<svg xmlns="http://www.w3.org/2000/svg" width="{size}" height="{size}" viewBox="0 0 64 64">
  <path d="M32 4L8 16v16c0 14 10 24 24 28 14-4 24-14 24-28V16L32 4z" fill="none" stroke="{color}" stroke-width="2.5"/>
  <polyline points="22,32 30,40 44,24" fill="none" stroke="{color}" stroke-width="2.5" stroke-linecap="round" stroke-linejoin="round"/>
</svg>'''


def svg_icon_target(size=64, color="#7E00FF"):
    return f'''<svg xmlns="http://www.w3.org/2000/svg" width="{size}" height="{size}" viewBox="0 0 64 64">
  <circle cx="32" cy="32" r="26" fill="none" stroke="{color}" stroke-width="2"/>
  <circle cx="32" cy="32" r="17" fill="none" stroke="{color}" stroke-width="2" opacity="0.7"/>
  <circle cx="32" cy="32" r="8" fill="none" stroke="{color}" stroke-width="2" opacity="0.5"/>
  <circle cx="32" cy="32" r="3" fill="{color}"/>
</svg>'''


def svg_icon_gear(size=64, color="#7E00FF"):
    return f'''<svg xmlns="http://www.w3.org/2000/svg" width="{size}" height="{size}" viewBox="0 0 64 64">
  <circle cx="32" cy="32" r="10" fill="none" stroke="{color}" stroke-width="2.5"/>
  <path d="M32 6v6m0 40v6m-26 6h6m40 0h6M11.5 11.5l4.2 4.2m32.6 32.6l4.2 4.2M11.5 52.5l4.2-4.2m32.6-32.6l4.2-4.2"
    fill="none" stroke="{color}" stroke-width="2.5" stroke-linecap="round" transform="translate(0,-6)"/>
</svg>'''


def svg_icon_chart_up(size=64, color="#7E00FF"):
    return f'''<svg xmlns="http://www.w3.org/2000/svg" width="{size}" height="{size}" viewBox="0 0 64 64">
  <polyline points="8,50 24,34 36,42 56,14" fill="none" stroke="{color}" stroke-width="2.5" stroke-linecap="round" stroke-linejoin="round"/>
  <polyline points="44,14 56,14 56,26" fill="none" stroke="{color}" stroke-width="2.5" stroke-linecap="round" stroke-linejoin="round"/>
</svg>'''


def svg_icon_brain(size=64, color="#7E00FF"):
    return f'''<svg xmlns="http://www.w3.org/2000/svg" width="{size}" height="{size}" viewBox="0 0 64 64">
  <path d="M32 56V12" fill="none" stroke="{color}" stroke-width="1.5" opacity="0.4"/>
  <path d="M22 16c-8 0-12 6-12 12s4 10 6 12c-2 2-4 6-2 10s6 6 10 6h8" fill="none" stroke="{color}" stroke-width="2.5" stroke-linecap="round"/>
  <path d="M42 16c8 0 12 6 12 12s-4 10-6 12c2 2 4 6 2 10s-6 6-10 6h-8" fill="none" stroke="{color}" stroke-width="2.5" stroke-linecap="round"/>
</svg>'''


def svg_icon_database(size=64, color="#7E00FF"):
    return f'''<svg xmlns="http://www.w3.org/2000/svg" width="{size}" height="{size}" viewBox="0 0 64 64">
  <ellipse cx="32" cy="16" rx="22" ry="8" fill="none" stroke="{color}" stroke-width="2"/>
  <path d="M10 16v32c0 4.4 9.8 8 22 8s22-3.6 22-8V16" fill="none" stroke="{color}" stroke-width="2"/>
  <path d="M10 32c0 4.4 9.8 8 22 8s22-3.6 22-8" fill="none" stroke="{color}" stroke-width="1.5" opacity="0.5"/>
</svg>'''


def svg_icon_refresh(size=64, color="#7E00FF"):
    return f'''<svg xmlns="http://www.w3.org/2000/svg" width="{size}" height="{size}" viewBox="0 0 64 64">
  <path d="M48 18A20 20 0 0 0 14 32" fill="none" stroke="{color}" stroke-width="2.5" stroke-linecap="round"/>
  <path d="M16 46A20 20 0 0 0 50 32" fill="none" stroke="{color}" stroke-width="2.5" stroke-linecap="round"/>
  <polyline points="48,10 48,20 38,20" fill="none" stroke="{color}" stroke-width="2.5" stroke-linecap="round" stroke-linejoin="round"/>
  <polyline points="16,54 16,44 26,44" fill="none" stroke="{color}" stroke-width="2.5" stroke-linecap="round" stroke-linejoin="round"/>
</svg>'''


def svg_icon_layers(size=64, color="#7E00FF"):
    return f'''<svg xmlns="http://www.w3.org/2000/svg" width="{size}" height="{size}" viewBox="0 0 64 64">
  <polygon points="32,8 58,24 32,40 6,24" fill="none" stroke="{color}" stroke-width="2" stroke-linejoin="round"/>
  <polyline points="6,32 32,48 58,32" fill="none" stroke="{color}" stroke-width="2" stroke-linejoin="round" opacity="0.6"/>
  <polyline points="6,40 32,56 58,40" fill="none" stroke="{color}" stroke-width="2" stroke-linejoin="round" opacity="0.35"/>
</svg>'''


def svg_icon_document(size=64, color="#7E00FF"):
    return f'''<svg xmlns="http://www.w3.org/2000/svg" width="{size}" height="{size}" viewBox="0 0 64 64">
  <path d="M16 4h22l14 14v38a4 4 0 0 1-4 4H16a4 4 0 0 1-4-4V8a4 4 0 0 1 4-4z" fill="none" stroke="{color}" stroke-width="2"/>
  <path d="M38 4v14h14" fill="none" stroke="{color}" stroke-width="2" stroke-linejoin="round"/>
  <line x1="20" y1="28" x2="44" y2="28" stroke="{color}" stroke-width="1.5" opacity="0.6"/>
  <line x1="20" y1="36" x2="44" y2="36" stroke="{color}" stroke-width="1.5" opacity="0.6"/>
  <line x1="20" y1="44" x2="36" y2="44" stroke="{color}" stroke-width="1.5" opacity="0.6"/>
</svg>'''


def svg_icon_check(size=64, color="#7E00FF"):
    return f'''<svg xmlns="http://www.w3.org/2000/svg" width="{size}" height="{size}" viewBox="0 0 64 64">
  <circle cx="32" cy="32" r="26" fill="none" stroke="{color}" stroke-width="2.5"/>
  <polyline points="20,32 28,42 44,22" fill="none" stroke="{color}" stroke-width="3" stroke-linecap="round" stroke-linejoin="round"/>
</svg>'''


def svg_icon_people(size=64, color="#7E00FF"):
    return f'''<svg xmlns="http://www.w3.org/2000/svg" width="{size}" height="{size}" viewBox="0 0 64 64">
  <circle cx="22" cy="20" r="8" fill="none" stroke="{color}" stroke-width="2"/>
  <path d="M6 50c0-10 8-16 16-16s16 6 16 16" fill="none" stroke="{color}" stroke-width="2" stroke-linecap="round"/>
  <circle cx="44" cy="20" r="8" fill="none" stroke="{color}" stroke-width="2" opacity="0.6"/>
  <path d="M34 50c0-10 6-16 14-16s10 6 10 16" fill="none" stroke="{color}" stroke-width="2" stroke-linecap="round" opacity="0.6"/>
</svg>'''


def svg_icon_lightbulb(size=64, color="#7E00FF"):
    return f'''<svg xmlns="http://www.w3.org/2000/svg" width="{size}" height="{size}" viewBox="0 0 64 64">
  <path d="M24 44v4a8 8 0 0 0 16 0v-4" fill="none" stroke="{color}" stroke-width="2"/>
  <path d="M32 8a18 18 0 0 0-12 32h24A18 18 0 0 0 32 8z" fill="none" stroke="{color}" stroke-width="2.5"/>
  <line x1="26" y1="52" x2="38" y2="52" stroke="{color}" stroke-width="1.5" opacity="0.5"/>
</svg>'''


# ══════════════════════════════════════════════════════════════
# Helper: add text paragraph
# ══════════════════════════════════════════════════════════════

def add_bullet(tf, text, bold=False, color_rgb=TEXT_BODY, size=14,
               space_before=6, prefix=""):
    p = tf.add_paragraph()
    p.space_before = Pt(space_before)
    if prefix:
        r = p.add_run()
        r.text = prefix
        r.font.size = Pt(size)
        r.font.bold = True
        r.font.color.rgb = PURPLE
        r.font.name = FONT
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color_rgb
    r.font.name = FONT
    return p


def setup_content_slide(breadcrumb, title, lead=""):
    """Create a content slide with breadcrumb, title, lead, footer, and GT."""
    slide = prs.slides.add_slide(prs.slide_layouts[2])
    # Breadcrumb
    for ph in slide.placeholders:
        idx = ph.placeholder_format.idx
        if idx == 11:
            ph.text = breadcrumb
            for p in ph.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(14)
                    r.font.bold = True
                    r.font.color.rgb = PURPLE
                    r.font.name = FONT
        elif idx == 0:
            ph.text = title
            for p in ph.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(28)
                    r.font.bold = True
                    r.font.color.rgb = BLACK
                    r.font.name = FONT
        elif idx == 10:
            if lead:
                ph.text = lead
                for p in ph.text_frame.paragraphs:
                    for r in p.runs:
                        r.font.size = Pt(18)
                        r.font.color.rgb = BLACK
                        r.font.name = FONT
            else:
                ph.text = ""
    set_footer(slide)
    add_gt_to_slide(slide)
    return slide


# ══════════════════════════════════════════════════════════════
# SLIDE 1 — Cover (Layout 0)
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[0])
for ph in slide.placeholders:
    idx = ph.placeholder_format.idx
    if idx == 0:
        ph.text = "Adoption-first scoring for\nenterprise AI feasibility"
        for p in ph.text_frame.paragraphs:
            for r in p.runs:
                r.font.size = Pt(40)
                r.font.bold = True
                r.font.color.rgb = BLACK
                r.font.name = FONT
    elif idx == 1:
        ph.text = "How to identify the skills that will actually get used"
        for p in ph.text_frame.paragraphs:
            for r in p.runs:
                r.font.size = Pt(18)
                r.font.color.rgb = BLACK
                r.font.name = FONT
    elif idx == 12:
        ph.text = "Aera Skill Feasibility Engine  |  Methodology Briefing"
        for p in ph.text_frame.paragraphs:
            for r in p.runs:
                r.font.size = Pt(14)
                r.font.color.rgb = BLACK
                r.font.name = FONT
    elif idx == 2:
        ph.text = "2026"
        for p in ph.text_frame.paragraphs:
            for r in p.runs:
                r.font.size = Pt(14)
                r.font.color.rgb = TEXT_SUB
                r.font.name = FONT
add_logo_to_cover(slide)


# ══════════════════════════════════════════════════════════════
# SLIDE 2 — The adoption failure problem (Pattern D)
# ══════════════════════════════════════════════════════════════
slide = setup_content_slide(
    "The Problem",
    "Most AI feasibility assessments score the wrong thing",
    "Traditional assessments weight technical capability highest. "
    "The result: technically sound features that nobody adopts."
)
col_w = (CW - 0.20) / 2
# Left card
lx = ML
card_h = AH - 0.10
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
    Inches(lx), Inches(CY), Inches(col_w), Inches(card_h))
shape.fill.solid()
shape.fill.fore_color.rgb = BG_LIGHT
shape.line.fill.background()
# Top bar
bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
    Inches(lx), Inches(CY), Inches(col_w), Inches(0.05))
bar.fill.solid()
bar.fill.fore_color.rgb = LIGHT_GRAY
bar.line.fill.background()

tb = slide.shapes.add_textbox(Inches(lx + 0.20), Inches(CY + 0.20),
                               Inches(col_w - 0.40), Inches(card_h - 0.30))
tf = tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
r = p.add_run()
r.text = "Traditional approach"
r.font.size = Pt(18)
r.font.bold = True
r.font.color.rgb = TEXT_SUB
r.font.name = FONT
add_bullet(tf, "Technical capability weighted 50%", space_before=14)
add_bullet(tf, "Adoption factors weighted 30%", space_before=8)
add_bullet(tf, "Value & ROI weighted 20%", space_before=8)
add_bullet(tf, "Outcome: 60%+ of enterprise AI features fail at adoption, not technology",
           space_before=16, bold=True, color_rgb=BLACK)

# Right card
rx = ML + col_w + 0.20
shape2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
    Inches(rx), Inches(CY), Inches(col_w), Inches(card_h))
shape2.fill.solid()
shape2.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
shape2.line.color.rgb = PURPLE
shape2.line.width = Pt(1.5)
# Top bar
bar2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
    Inches(rx), Inches(CY), Inches(col_w), Inches(0.05))
bar2.fill.solid()
bar2.fill.fore_color.rgb = PURPLE
bar2.line.fill.background()

tb2 = slide.shapes.add_textbox(Inches(rx + 0.20), Inches(CY + 0.20),
                                Inches(col_w - 0.40), Inches(card_h - 0.30))
tf2 = tb2.text_frame
tf2.word_wrap = True
p = tf2.paragraphs[0]
r = p.add_run()
r.text = "Adoption-first approach"
r.font.size = Pt(18)
r.font.bold = True
r.font.color.rgb = PURPLE
r.font.name = FONT
add_bullet(tf2, "Technical capability weighted 30%", space_before=14)
add_bullet(tf2, "Adoption realism weighted 45%", space_before=8, bold=True, color_rgb=PURPLE)
add_bullet(tf2, "Value & ROI weighted 25%", space_before=8)
add_bullet(tf2, "Outcome: every evaluation forces the question \"will anyone actually use this?\"",
           space_before=16, bold=True, color_rgb=BLACK)

add_svg_native(slide, svg_icon_target(48, "#D8D8D8"),
               x=lx + col_w - 0.65, y=CY + card_h - 0.70, w=0.45, h=0.45)
add_svg_native(slide, svg_icon_target(48, "#7E00FF"),
               x=rx + col_w - 0.65, y=CY + card_h - 0.70, w=0.45, h=0.45)


# ══════════════════════════════════════════════════════════════
# SLIDE 3 — What if we scored differently? (Pattern E — 3 cols)
# ══════════════════════════════════════════════════════════════
slide = setup_content_slide(
    "The Problem",
    "What the engine delivers in a single overnight run",
    "One CLI command. Fully unattended batch processing. "
    "Actionable output an SE team can execute on."
)
col3_w = (CW - 2 * 0.15) / 3
items3 = [
    ("Three-lens scoring", [
        "Adoption realism at 45%",
        "9 sub-dimensions scored 0-3",
        "Composite threshold gate at 0.60",
    ], svg_icon_brain(48)),
    ("Automated triage", [
        "5 data-driven red flags",
        "3-tier priority binning",
        "Archetype classification",
    ], svg_icon_shield(48)),
    ("Grounded simulation", [
        "Decision flow diagrams",
        "Component maps (YAML)",
        "Mock tests + integration specs",
    ], svg_icon_document(48)),
]
for i, (title_txt, bullets, icon_svg) in enumerate(items3):
    cx = ML + i * (col3_w + 0.15)
    card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(cx), Inches(CY), Inches(col3_w), Inches(AH))
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = LIGHT_GRAY
    card.line.width = Pt(0.75)
    # Top bar
    tb_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(cx), Inches(CY), Inches(col3_w), Inches(0.05))
    tb_bar.fill.solid()
    tb_bar.fill.fore_color.rgb = PURPLE
    tb_bar.line.fill.background()
    # Icon
    add_svg_native(slide, icon_svg,
                   x=cx + (col3_w - 0.40) / 2, y=CY + 0.20, w=0.40, h=0.40)
    # Title
    tt = slide.shapes.add_textbox(Inches(cx + 0.12), Inches(CY + 0.70),
                                   Inches(col3_w - 0.24), Inches(0.40))
    ttf = tt.text_frame
    ttf.word_wrap = True
    p = ttf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = title_txt
    r.font.size = Pt(16)
    r.font.bold = True
    r.font.color.rgb = PURPLE
    r.font.name = FONT
    # Bullets
    bt = slide.shapes.add_textbox(Inches(cx + 0.15), Inches(CY + 1.20),
                                   Inches(col3_w - 0.30), Inches(AH - 1.30))
    btf = bt.text_frame
    btf.word_wrap = True
    for j, b in enumerate(bullets):
        p = btf.add_paragraph() if j > 0 else btf.paragraphs[0]
        p.space_before = Pt(8)
        r = p.add_run()
        r.text = f"  {b}"
        r.font.size = Pt(14)
        r.font.color.rgb = TEXT_BODY
        r.font.name = FONT


# ══════════════════════════════════════════════════════════════
# SLIDE 4 — Three lenses, adoption-weighted (Pattern C — table + chart)
# ══════════════════════════════════════════════════════════════
slide = setup_content_slide(
    "Three-Lens Scoring",
    "Adoption realism gets the highest weight at 45%",
    "Three lenses produce a composite score. Traditional assessments "
    "weight technical capability first. This engine deliberately inverts that."
)
# Grouped bar chart
chart_svg = svg_grouped_bar_chart(
    categories=["Technical Feasibility", "Adoption Realism", "Value &amp; Efficiency"],
    series1=[50, 30, 20],
    series2=[30, 45, 25],
    label1="Traditional",
    label2="Adoption-First",
    w=700, h=280
)
add_svg_native(slide, chart_svg, x=ML + 0.15, y=CY + 0.05, w=7.80, h=3.60)

# Right-side summary card
rx = ML + 8.20
rw = CW - 8.20
card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
    Inches(rx), Inches(CY), Inches(rw), Inches(AH))
card.fill.solid()
card.fill.fore_color.rgb = BG_LIGHT
card.line.fill.background()
tb = slide.shapes.add_textbox(Inches(rx + 0.20), Inches(CY + 0.15),
                               Inches(rw - 0.40), Inches(AH - 0.30))
tf = tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
r = p.add_run()
r.text = "Weight rationale"
r.font.size = Pt(16)
r.font.bold = True
r.font.color.rgb = DEEP_PURPLE
r.font.name = FONT

lens_items = [
    ("Technical (30%)", "Can we build it? Necessary but not dominant."),
    ("Adoption (45%)", "Will anyone use it? The differentiator."),
    ("Value (25%)", "Is it worth it? ROI and simulation viability."),
]
for label, desc in lens_items:
    p = tf.add_paragraph()
    p.space_before = Pt(14)
    r = p.add_run()
    r.text = label
    r.font.size = Pt(14)
    r.font.bold = True
    r.font.color.rgb = PURPLE
    r.font.name = FONT
    p2 = tf.add_paragraph()
    p2.space_before = Pt(4)
    r2 = p2.add_run()
    r2.text = desc
    r2.font.size = Pt(14)
    r2.font.color.rgb = TEXT_BODY
    r2.font.name = FONT


# ══════════════════════════════════════════════════════════════
# SLIDE 5 — Technical Feasibility lens (Pattern M)
# ══════════════════════════════════════════════════════════════
slide = setup_content_slide(
    "Three-Lens Scoring",
    "Technical feasibility: necessary but deliberately not dominant",
    "Three sub-dimensions scored 0-3. Max raw score: 9. Weight: 30%."
)
LABEL_W = 2.0
subdims = [
    ("Data\nReadiness", "Are the needed data sources available and mappable to Aera's ingestion layer? "
     "Evaluates whether source systems (ERP, CRM, operational databases) can feed the required Streams."),
    ("Aera Platform\nFit", "Does the opportunity's archetype (DETERMINISTIC, AGENTIC, or GENERATIVE) "
     "map cleanly to an Aera component? Process Builder for rules, Agent Teams for orchestration, "
     "Cortex for ML models."),
    ("Archetype\nConfidence", "How certain is the archetype classification? Directly exported "
     "classifications score higher than heuristically inferred ones."),
]
row_h = (AH - 0.20) / 3
for i, (label, desc) in enumerate(subdims):
    ry = CY + i * (row_h + 0.10)
    lbl = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(ML), Inches(ry), Inches(LABEL_W), Inches(row_h))
    lbl.fill.solid()
    lbl.fill.fore_color.rgb = DEEP_PURPLE
    lbl.line.fill.background()
    lt = slide.shapes.add_textbox(Inches(ML), Inches(ry),
                                   Inches(LABEL_W), Inches(row_h))
    ltf = lt.text_frame
    ltf.word_wrap = True
    ltf.margin_top = Inches(row_h * 0.25)
    p = ltf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = label
    r.font.size = Pt(14)
    r.font.bold = True
    r.font.color.rgb = WHITE
    r.font.name = FONT

    ct = slide.shapes.add_textbox(Inches(ML + LABEL_W + 0.12), Inches(ry + 0.10),
                                   Inches(CW - LABEL_W - 0.12), Inches(row_h - 0.15))
    ctf = ct.text_frame
    ctf.word_wrap = True
    p = ctf.paragraphs[0]
    r = p.add_run()
    r.text = desc
    r.font.size = Pt(14)
    r.font.color.rgb = TEXT_BODY
    r.font.name = FONT


# ══════════════════════════════════════════════════════════════
# SLIDE 6 — Adoption Realism lens (Pattern J)
# ══════════════════════════════════════════════════════════════
slide = setup_content_slide(
    "Three-Lens Scoring",
    "Adoption realism at 45% prevents operationally dead features",
    "Four sub-dimensions scored 0-3. Max raw score: 12. Weight: 45%."
)
adoption_items = [
    ("Decision\nDensity", "Are there actual, recurring business decisions to automate? "
     "Opportunities with zero decision density are flagged as DEAD_ZONE."),
    ("Financial\nGravity", "Does the outcome have measurable financial impact? "
     "Evaluates revenue effect, cost reduction potential, and financial rating distribution."),
    ("Impact\nProximity", "Is the benefit felt immediately by users, or is it distant and abstract? "
     "First-order impact (direct user value) scores higher than second-order (indirect)."),
    ("Confidence\nSignal", "Does the organization have conviction about this opportunity? "
     "Derived from rating confidence across constituent activities."),
]
row_h4 = (AH - 0.75 - 0.30) / 4
for i, (label, desc) in enumerate(adoption_items):
    ry = CY + i * (row_h4 + 0.10)
    lbl = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(ML), Inches(ry), Inches(1.60), Inches(row_h4))
    lbl.fill.solid()
    lbl.fill.fore_color.rgb = PURPLE if i % 2 == 0 else DEEP_PURPLE
    lbl.line.fill.background()
    lt = slide.shapes.add_textbox(Inches(ML), Inches(ry),
                                   Inches(1.60), Inches(row_h4))
    ltf = lt.text_frame
    ltf.word_wrap = True
    ltf.margin_top = Inches(row_h4 * 0.20)
    p = ltf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = label
    r.font.size = Pt(14)
    r.font.bold = True
    r.font.color.rgb = WHITE
    r.font.name = FONT

    ct = slide.shapes.add_textbox(Inches(ML + 1.72), Inches(ry + 0.06),
                                   Inches(CW - 1.72), Inches(row_h4 - 0.10))
    ctf = ct.text_frame
    ctf.word_wrap = True
    p = ctf.paragraphs[0]
    r = p.add_run()
    r.text = desc
    r.font.size = Pt(14)
    r.font.color.rgb = TEXT_BODY
    r.font.name = FONT

# Highlight bar
hl_y = CY + 4 * (row_h4 + 0.10) + 0.05
hl_card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
    Inches(ML), Inches(hl_y), Inches(CW), Inches(0.65))
hl_card.fill.solid()
hl_card.fill.fore_color.rgb = LP_BG
hl_card.line.color.rgb = PURPLE
hl_card.line.width = Pt(1)
hl_tb = slide.shapes.add_textbox(Inches(ML + 0.30), Inches(hl_y + 0.05),
                                  Inches(CW - 0.60), Inches(0.55))
hl_tf = hl_tb.text_frame
hl_tf.word_wrap = True
p = hl_tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = "45% weight prevents technically perfect but operationally dead implementations"
r.font.size = Pt(16)
r.font.bold = True
r.font.color.rgb = DEEP_PURPLE
r.font.name = FONT


# ══════════════════════════════════════════════════════════════
# SLIDE 7 — Value & Efficiency lens (Pattern R)
# ══════════════════════════════════════════════════════════════
slide = setup_content_slide(
    "Three-Lens Scoring",
    "Value and efficiency: grounded in real client financials",
    "Two sub-dimensions scored 0-3. Max raw score: 6. Weight: 25%."
)
value_items = [
    ("Value\nDensity",
     "Revenue impact as a percentage of company annual revenue. "
     "The engine ingests actual financial data from the hierarchy export "
     "(annual revenue, COGS, operating costs) to compute value density. "
     "This prevents inflated ROI claims detached from real business scale."),
    ("Simulation\nViability",
     "Can the impact be quantified and tested before full rollout? "
     "Evaluates whether the opportunity has enough structured data "
     "and decision patterns to generate a credible mock decision test "
     "with sample inputs, expected outputs, and integration points."),
]
val_row_h = (AH - 0.10) / 2
for i, (label, desc) in enumerate(value_items):
    ry = CY + i * (val_row_h + 0.10)
    lbl = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(ML), Inches(ry), Inches(2.0), Inches(val_row_h))
    lbl.fill.solid()
    lbl.fill.fore_color.rgb = DEEP_PURPLE
    lbl.line.fill.background()
    lt = slide.shapes.add_textbox(Inches(ML), Inches(ry),
                                   Inches(2.0), Inches(val_row_h))
    ltf = lt.text_frame
    ltf.word_wrap = True
    ltf.margin_top = Inches(val_row_h * 0.30)
    p = ltf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = label
    r.font.size = Pt(14)
    r.font.bold = True
    r.font.color.rgb = WHITE
    r.font.name = FONT

    ct = slide.shapes.add_textbox(Inches(ML + 2.12), Inches(ry + 0.12),
                                   Inches(CW - 2.12), Inches(val_row_h - 0.20))
    ctf = ct.text_frame
    ctf.word_wrap = True
    p = ctf.paragraphs[0]
    r = p.add_run()
    r.text = desc
    r.font.size = Pt(14)
    r.font.color.rgb = TEXT_BODY
    r.font.name = FONT


# ══════════════════════════════════════════════════════════════
# SLIDE 8 — Composite scoring & 0.60 gate (Pattern I — chevron + detail)
# ══════════════════════════════════════════════════════════════
slide = setup_content_slide(
    "Three-Lens Scoring",
    "Composite score gate at 0.60 focuses simulation effort",
    "Each lens is normalized 0-1, weighted, and summed. Only opportunities at or above 0.60 advance."
)
# Chevron flow
chev_svg = svg_chevron_flow(["Normalize\n0-1", "Apply\nWeights", "Sum\nComposite", "Gate\n>= 0.60"],
                             w=800, h=60)
add_svg_native(slide, chev_svg, x=ML, y=CY + 0.05, w=CW, h=0.65)

# Worked example card
ex_y = CY + 0.85
ex_h = AH - 0.85
card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
    Inches(ML), Inches(ex_y), Inches(CW), Inches(ex_h))
card.fill.solid()
card.fill.fore_color.rgb = BG_LIGHT
card.line.fill.background()

tb = slide.shapes.add_textbox(Inches(ML + 0.20), Inches(ex_y + 0.10),
                               Inches(CW - 0.40), Inches(0.35))
tf = tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
r = p.add_run()
r.text = "Worked example: Supply Chain Demand Sensing"
r.font.size = Pt(16)
r.font.bold = True
r.font.color.rgb = DEEP_PURPLE
r.font.name = FONT

# Stacked bar showing composition
bar_svg = svg_stacked_horiz_bar(
    segments=[30, 45, 25],
    labels=["Tech 0.78 x 0.30", "Adopt 0.83 x 0.45", "Value 0.67 x 0.25"],
    total_w=700, h=45,
    colors=["#5F0095", "#7E00FF", "#D8D8D8"]
)
add_svg_native(slide, bar_svg, x=ML + 0.20, y=ex_y + 0.55, w=8.00, h=0.55)

# Calculation text
calc_tb = slide.shapes.add_textbox(Inches(ML + 0.20), Inches(ex_y + 1.25),
                                    Inches(CW - 3.80), Inches(ex_h - 1.45))
calc_tf = calc_tb.text_frame
calc_tf.word_wrap = True
items_calc = [
    ("Technical:", "7/9 = 0.78 normalized x 0.30 weight = 0.233"),
    ("Adoption:", "10/12 = 0.83 normalized x 0.45 weight = 0.375"),
    ("Value:", "4/6 = 0.67 normalized x 0.25 weight = 0.167"),
    ("Composite:", "0.233 + 0.375 + 0.167 = 0.775  (above 0.60 gate)"),
]
for j, (lbl, val) in enumerate(items_calc):
    p = calc_tf.add_paragraph() if j > 0 else calc_tf.paragraphs[0]
    p.space_before = Pt(10)
    rl = p.add_run()
    rl.text = f"  {lbl}  "
    rl.font.size = Pt(14)
    rl.font.bold = True
    rl.font.color.rgb = PURPLE
    rl.font.name = FONT
    rv = p.add_run()
    rv.text = val
    rv.font.size = Pt(14)
    rv.font.color.rgb = TEXT_BODY
    rv.font.name = FONT

# Result badge — placed just below the card to avoid overlap
result_y = ex_y + ex_h - 0.48
res_card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
    Inches(ML + CW - 3.20), Inches(result_y), Inches(3.00), Inches(0.38))
res_card.fill.solid()
res_card.fill.fore_color.rgb = LP_BG
res_card.line.color.rgb = PURPLE
res_card.line.width = Pt(1)
res_tb = slide.shapes.add_textbox(Inches(ML + CW - 3.15), Inches(result_y + 0.02),
                                   Inches(2.90), Inches(0.34))
res_tf = res_tb.text_frame
p = res_tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = "PROMOTED TO SIMULATION"
r.font.size = Pt(14)
r.font.bold = True
r.font.color.rgb = DEEP_PURPLE
r.font.name = FONT


# ══════════════════════════════════════════════════════════════
# SLIDE 9 — Red flags (Pattern C — table)
# ══════════════════════════════════════════════════════════════
slide = setup_content_slide(
    "Automated Triage",
    "Five red flags catch data problems before scoring begins",
    "Before any scoring, the engine detects structural problems in the hierarchy data. "
    "Each flag triggers a deterministic action: skip, demote, or flag for review."
)
# SVG action matrix
matrix_svg = svg_action_matrix(w=750, h=220)
add_svg_native(slide, matrix_svg, x=ML + 0.15, y=CY + 0.10, w=CW - 0.30, h=3.10)

# Explanation below
note_tb = slide.shapes.add_textbox(Inches(ML + 0.15), Inches(CY + 3.35),
                                    Inches(CW - 0.30), Inches(AH - 3.50))
note_tf = note_tb.text_frame
note_tf.word_wrap = True
actions_desc = [
    ("SKIP", "removes the opportunity from further processing entirely"),
    ("DEMOTE", "drops the opportunity one tier (e.g., Tier 1 to Tier 2)"),
    ("FLAG", "marks the issue for review but continues processing"),
]
for j, (act, desc) in enumerate(actions_desc):
    p = note_tf.add_paragraph() if j > 0 else note_tf.paragraphs[0]
    p.space_before = Pt(10)
    ra = p.add_run()
    ra.text = f"  {act}:  "
    ra.font.size = Pt(14)
    ra.font.bold = True
    ra.font.color.rgb = PURPLE
    ra.font.name = FONT
    rd = p.add_run()
    rd.text = desc
    rd.font.size = Pt(14)
    rd.font.color.rgb = TEXT_BODY
    rd.font.name = FONT


# ══════════════════════════════════════════════════════════════
# SLIDE 10 — Tier assignment (Pattern E — 3 cols)
# ══════════════════════════════════════════════════════════════
slide = setup_content_slide(
    "Automated Triage",
    "Three tiers bin opportunities by ROI potential",
    "Tier 1 is checked first. Red flags can demote. Tiers determine "
    "analysis depth and simulation priority."
)
tiers = [
    ("Tier 1: Premium", "#5F0095", [
        "quick_win flag = true",
        "Combined value > $5M",
        "Maximum ROI priority",
        "Full simulation package",
    ]),
    ("Tier 2: Moderate", "#7E00FF", [
        ">= 50% HIGH ai_suitability",
        "Good technical fit",
        "Standard simulation",
        "Scorecard + component map",
    ]),
    ("Tier 3: Default", "#D8D8D8", [
        "Everything else",
        "Baseline evaluation",
        "Scorecard only",
        "Triage summary in TSV",
    ]),
]
tier_col_w = (CW - 2 * 0.15) / 3
for i, (title_t, col_c, bullets_t) in enumerate(tiers):
    cx = ML + i * (tier_col_w + 0.15)
    card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(cx), Inches(CY), Inches(tier_col_w), Inches(AH))
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = LIGHT_GRAY
    card.line.width = Pt(0.75)
    # Top bar
    tbar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(cx), Inches(CY), Inches(tier_col_w), Inches(0.05))
    tbar.fill.solid()
    tbar.fill.fore_color.rgb = RGBColor.from_string(col_c[1:])
    tbar.line.fill.background()
    # Badge
    add_svg_native(slide, svg_callout_badge(str(i+1), 44, col_c),
                   x=cx + (tier_col_w - 0.40) / 2, y=CY + 0.20, w=0.40, h=0.40)
    # Title
    tt = slide.shapes.add_textbox(Inches(cx + 0.10), Inches(CY + 0.70),
                                   Inches(tier_col_w - 0.20), Inches(0.40))
    ttf = tt.text_frame
    ttf.word_wrap = True
    p = ttf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = title_t
    r.font.size = Pt(16)
    r.font.bold = True
    r.font.color.rgb = RGBColor.from_string(col_c[1:]) if col_c != "#D8D8D8" else TEXT_BODY
    r.font.name = FONT
    # Bullets
    bt = slide.shapes.add_textbox(Inches(cx + 0.15), Inches(CY + 1.20),
                                   Inches(tier_col_w - 0.30), Inches(AH - 1.30))
    btf = bt.text_frame
    btf.word_wrap = True
    for j, b in enumerate(bullets_t):
        p = btf.add_paragraph() if j > 0 else btf.paragraphs[0]
        p.space_before = Pt(10)
        r = p.add_run()
        r.text = f"  {b}"
        r.font.size = Pt(14)
        r.font.color.rgb = TEXT_BODY
        r.font.name = FONT


# ══════════════════════════════════════════════════════════════
# SLIDE 11 — Archetype routing (Pattern B — 2x2 becomes 3-card row)
# ══════════════════════════════════════════════════════════════
slide = setup_content_slide(
    "Automated Triage",
    "Archetype routing maps opportunities to platform components",
    "Each opportunity is classified as DETERMINISTIC, AGENTIC, or GENERATIVE. "
    "This determines which Aera components are used in simulation."
)
archetypes = [
    ("DETERMINISTIC", "~56%", [
        "Process Builder",
        "IF/Rules/Transaction nodes",
        "Structured workflows",
        "Highest prevalence",
    ]),
    ("AGENTIC", "~43%", [
        "Agent Teams + Functions",
        "Multi-step orchestration",
        "Decision-density driven",
        "Growing category",
    ]),
    ("GENERATIVE", "<1%", [
        "Agent Teams + LLM Agents",
        "Cortex model integration",
        "Unstructured reasoning",
        "Rare in enterprise ops",
    ]),
]
arch_w = (CW - 2 * 0.15) / 3
for i, (aname, pct, bullets_a) in enumerate(archetypes):
    cx = ML + i * (arch_w + 0.15)
    # Card
    card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(cx), Inches(CY), Inches(arch_w), Inches(AH))
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = LIGHT_GRAY
    card.line.width = Pt(0.75)
    # Header area
    hdr = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(cx), Inches(CY), Inches(arch_w), Inches(0.90))
    hdr.fill.solid()
    hdr.fill.fore_color.rgb = DEEP_PURPLE if i == 0 else (PURPLE if i == 1 else LIGHT_GRAY)
    hdr.line.fill.background()
    ht = slide.shapes.add_textbox(Inches(cx + 0.10), Inches(CY + 0.08),
                                   Inches(arch_w - 0.20), Inches(0.44))
    htf = ht.text_frame
    htf.word_wrap = True
    p = htf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = aname
    r.font.size = Pt(16)
    r.font.bold = True
    r.font.color.rgb = WHITE if i < 2 else TEXT_BODY
    r.font.name = FONT
    # Percentage
    pt_t = slide.shapes.add_textbox(Inches(cx + 0.10), Inches(CY + 0.48),
                                     Inches(arch_w - 0.20), Inches(0.36))
    ptf = pt_t.text_frame
    p = ptf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = pct
    r.font.size = Pt(20)
    r.font.bold = True
    r.font.color.rgb = WHITE if i < 2 else TEXT_BODY
    r.font.name = FONT
    # Bullets
    bt = slide.shapes.add_textbox(Inches(cx + 0.15), Inches(CY + 1.10),
                                   Inches(arch_w - 0.30), Inches(AH - 1.20))
    btf = bt.text_frame
    btf.word_wrap = True
    for j, b in enumerate(bullets_a):
        p = btf.add_paragraph() if j > 0 else btf.paragraphs[0]
        p.space_before = Pt(10)
        r = p.add_run()
        r.text = f"  {b}"
        r.font.size = Pt(14)
        r.font.color.rgb = TEXT_BODY
        r.font.name = FONT


# ══════════════════════════════════════════════════════════════
# SLIDE 12 — Grounded simulation (Pattern V — chevron + cards)
# ══════════════════════════════════════════════════════════════
slide = setup_content_slide(
    "Simulation",
    "Every simulation artifact is grounded in a real knowledge base",
    "The engine bundles 21 UI components, 22 Process Builder nodes, and 7 workflow patterns. "
    "Every reference is validated. Confidence: confirmed (KB match) or inferred (flagged)."
)
def _chevron_band_svg(labels, w=800, h=55):
    n = len(labels)
    sw = w / n
    td = h * 0.38
    alts = ["#5F0095", "#7E00FF"]
    polys = []
    for i, lbl in enumerate(labels):
        c = alts[i % 2]
        x0 = i * sw
        x1 = x0 + sw
        mid = h / 2
        pts = f"{x0},0 {x1-td},0 {x1},{mid} {x1-td},{h} {x0},{h}"
        polys.append(
            f'<polygon points="{pts}" fill="{c}"/>'
            f'<text x="{x0+sw/2}" y="{mid}" text-anchor="middle" dominant-baseline="middle" '
            f'fill="white" font-size="13" font-weight="bold" '
            f'font-family="Graphik, sans-serif">{lbl}</text>'
        )
    return f'<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 {w} {h}">{"".join(polys)}</svg>'

col_labels = ["Decision Flow", "Component Map", "Mock Test", "Integration"]
col_bullets = [
    ["Mermaid diagram", "Trigger to outcome", "Real PB node labels", "Visual workflow"],
    ["YAML structure", "Streams + Cortex", "PB nodes + Agents", "UI components"],
    ["Sample scenario", "Client financials", "Expected output", "Integration points"],
    ["Data flow topology", "Source systems", "Aera ingestion", "Processing to UI"],
]
col_icons = [svg_icon_layers(48), svg_icon_database(48), svg_icon_check(48), svg_icon_gear(48)]

chev_band = _chevron_band_svg(col_labels, w=800, h=50)
add_svg_native(slide, chev_band, x=ML, y=CY, w=CW, h=0.55)

N = len(col_labels)
card_w = (CW - (N - 1) * 0.10) / N
card_y = CY + 0.55 + 0.12
card_h = AH - 0.55 - 0.12
for i in range(N):
    cx = ML + i * (card_w + 0.10)
    card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(cx), Inches(card_y), Inches(card_w), Inches(card_h))
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = LIGHT_GRAY
    card.line.width = Pt(0.75)
    # Icon
    add_svg_native(slide, col_icons[i],
                   x=cx + (card_w - 0.35) / 2, y=card_y + 0.12, w=0.35, h=0.35)
    # Bullets
    bt = slide.shapes.add_textbox(Inches(cx + 0.10), Inches(card_y + 0.55),
                                   Inches(card_w - 0.20), Inches(card_h - 0.65))
    btf = bt.text_frame
    btf.word_wrap = True
    for j, b in enumerate(col_bullets[i]):
        p = btf.add_paragraph() if j > 0 else btf.paragraphs[0]
        p.space_before = Pt(6)
        r = p.add_run()
        r.text = f"  {b}"
        r.font.size = Pt(14)
        r.font.color.rgb = TEXT_BODY
        r.font.name = FONT


# ══════════════════════════════════════════════════════════════
# SLIDE 13 — Four simulation artifacts (Pattern E)
# ══════════════════════════════════════════════════════════════
slide = setup_content_slide(
    "Simulation",
    "Four artifacts per qualifying opportunity, all validated",
    "Each artifact is generated by a dedicated generator module. "
    "Knowledge base validation flags any component not in the bundled catalog."
)
artifacts = [
    ("Mermaid decision flow", [
        "Visual workflow: trigger to decision nodes to outcome",
        "Nodes labeled with real Process Builder types",
        "Validates Mermaid syntax automatically",
    ]),
    ("YAML component map", [
        "Which Streams, Cortex models, PB nodes, Agent Teams, and UI components",
        "Every component checked against 21 UI + 22 PB node catalog",
        "Confidence: confirmed (exact KB match) or inferred",
    ]),
    ("Mock decision test", [
        "Sample scenario using real client financial data",
        "Expected output with decision rationale",
        "Integration points for downstream systems",
    ]),
    ("Integration surface", [
        "Full data flow topology from source to UI",
        "Source systems, Aera ingestion, processing layer",
        "Identifies external system dependencies",
    ]),
]
art_col_w = (CW - 3 * 0.12) / 4
for i, (atitle, abullets) in enumerate(artifacts):
    cx = ML + i * (art_col_w + 0.12)
    card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(cx), Inches(CY), Inches(art_col_w), Inches(AH))
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = LIGHT_GRAY
    card.line.width = Pt(0.75)
    # Top bar
    tbar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(cx), Inches(CY), Inches(art_col_w), Inches(0.05))
    tbar.fill.solid()
    tbar.fill.fore_color.rgb = PURPLE
    tbar.line.fill.background()
    # Badge
    add_svg_native(slide, svg_callout_badge(str(i+1), 36),
                   x=cx + (art_col_w - 0.32) / 2, y=CY + 0.15, w=0.32, h=0.32)
    # Title
    tt = slide.shapes.add_textbox(Inches(cx + 0.08), Inches(CY + 0.55),
                                   Inches(art_col_w - 0.16), Inches(0.50))
    ttf = tt.text_frame
    ttf.word_wrap = True
    p = ttf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = atitle
    r.font.size = Pt(14)
    r.font.bold = True
    r.font.color.rgb = PURPLE
    r.font.name = FONT
    # Bullets
    bt = slide.shapes.add_textbox(Inches(cx + 0.08), Inches(CY + 1.10),
                                   Inches(art_col_w - 0.16), Inches(AH - 1.20))
    btf = bt.text_frame
    btf.word_wrap = True
    for j, b in enumerate(abullets):
        p = btf.add_paragraph() if j > 0 else btf.paragraphs[0]
        p.space_before = Pt(8)
        r = p.add_run()
        r.text = f"  {b}"
        r.font.size = Pt(14)
        r.font.color.rgb = TEXT_BODY
        r.font.name = FONT


# ══════════════════════════════════════════════════════════════
# SLIDE 14 — From score to spec (Pattern F — numbered summary)
# ══════════════════════════════════════════════════════════════
slide = setup_content_slide(
    "Simulation",
    "The evaluation directory an SE team can act on Monday morning",
    "Every run produces a structured output directory. "
    "All artifacts are machine-readable and human-reviewable."
)
deliverables = [
    ("Triage summary (TSV)", "Sortable spreadsheet of all opportunities with flags, tiers, and actions"),
    ("Feasibility scorecards", "9-dimension numeric grading for every opportunity across three lenses"),
    ("Tier 1 deep analysis", "Detailed report for premium opportunities with scoring reasoning"),
    ("Simulation packages", "Per-opportunity folders with flows, maps, tests, and integration specs"),
    ("Catalog analysis", "Meta-reflection on domain strengths, archetype distribution, and KB coverage"),
]
item_h = (AH - 0.10) / len(deliverables)
for i, (dtitle, ddesc) in enumerate(deliverables):
    iy = CY + i * item_h
    # Badge
    add_svg_native(slide, svg_callout_badge(str(i+1), 36),
                   x=ML + 0.05, y=iy + (item_h - 0.30) / 2, w=0.30, h=0.30)
    # Card
    card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(ML + 0.45), Inches(iy + 0.03), Inches(CW - 0.45), Inches(item_h - 0.06))
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = LIGHT_GRAY
    card.line.width = Pt(0.75)
    # Purple accent bar
    abar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(ML + 0.45), Inches(iy + 0.03), Inches(0.06), Inches(item_h - 0.06))
    abar.fill.solid()
    abar.fill.fore_color.rgb = PURPLE
    abar.line.fill.background()
    # Text
    tb = slide.shapes.add_textbox(Inches(ML + 0.65), Inches(iy + 0.06),
                                   Inches(CW - 0.80), Inches(item_h - 0.12))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = dtitle
    r.font.size = Pt(14)
    r.font.bold = True
    r.font.color.rgb = DEEP_PURPLE
    r.font.name = FONT
    p2 = tf.add_paragraph()
    p2.space_before = Pt(4)
    r2 = p2.add_run()
    r2.text = ddesc
    r2.font.size = Pt(14)
    r2.font.color.rgb = TEXT_BODY
    r2.font.name = FONT


# ══════════════════════════════════════════════════════════════
# SLIDE 15 — Overnight batch resilience (Pattern D — two-col)
# ══════════════════════════════════════════════════════════════
slide = setup_content_slide(
    "Resilience & Scale",
    "The engine runs overnight with zero babysitting",
    "Four resilience mechanisms ensure unattended operation completes reliably."
)
col_w15 = (CW - 0.20) / 2
# Left: resilience features
lcard = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
    Inches(ML), Inches(CY), Inches(col_w15), Inches(AH))
lcard.fill.solid()
lcard.fill.fore_color.rgb = WHITE
lcard.line.color.rgb = LIGHT_GRAY
lcard.line.width = Pt(0.75)
ltbar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
    Inches(ML), Inches(CY), Inches(col_w15), Inches(0.05))
ltbar.fill.solid()
ltbar.fill.fore_color.rgb = PURPLE
ltbar.line.fill.background()

resilience_items = [
    ("Checkpoint recovery", "Crash mid-run? Resume from last completed opportunity. No re-scoring."),
    ("Three-tier retry", "LLM failure triggers retry, then fallback prompt, then skip-and-log."),
    ("Git auto-commit", "Artifacts committed after each phase with markers for traceability."),
    ("Conservative confidence", "Overall confidence = MIN of lens confidences. Never overstates certainty."),
]
ltb = slide.shapes.add_textbox(Inches(ML + 0.15), Inches(CY + 0.15),
                                Inches(col_w15 - 0.30), Inches(AH - 0.20))
ltf = ltb.text_frame
ltf.word_wrap = True
for j, (rtitle, rdesc) in enumerate(resilience_items):
    if j > 0:
        p = ltf.add_paragraph()
        p.space_before = Pt(14)
    else:
        p = ltf.paragraphs[0]
        p.space_before = Pt(0)
    r = p.add_run()
    r.text = rtitle
    r.font.size = Pt(14)
    r.font.bold = True
    r.font.color.rgb = PURPLE
    r.font.name = FONT
    p2 = ltf.add_paragraph()
    p2.space_before = Pt(4)
    r2 = p2.add_run()
    r2.text = rdesc
    r2.font.size = Pt(14)
    r2.font.color.rgb = TEXT_BODY
    r2.font.name = FONT

# Right: architecture diagram
rcard = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
    Inches(ML + col_w15 + 0.20), Inches(CY), Inches(col_w15), Inches(AH))
rcard.fill.solid()
rcard.fill.fore_color.rgb = BG_LIGHT
rcard.line.fill.background()

layers_svg = svg_layers_diagram(
    ["CLI Command", "Checkpoint Gate", "LLM Scoring (retry/fallback)", "Simulation Gen", "Git Auto-Commit"],
    w=500, h=260
)
add_svg_native(slide, layers_svg,
               x=ML + col_w15 + 0.50, y=CY + 0.60, w=col_w15 - 0.60, h=3.60)


# ══════════════════════════════════════════════════════════════
# SLIDE 16 — Cloud acceleration (Pattern G — chart + sidebar)
# ══════════════════════════════════════════════════════════════
slide = setup_content_slide(
    "Resilience & Scale",
    "Optional cloud path reduces 17 hours to under 30 minutes",
    "Local Ollama path works fully offline. Optional vLLM on H100 enables "
    "concurrent scoring. Same logic, same artifacts, same quality."
)
# Timeline chart
timeline_svg = svg_timeline_compare(w=750, h=200)
add_svg_native(slide, timeline_svg, x=ML + 0.10, y=CY + 0.10, w=8.50, h=2.80)

# Side panel
px = ML + 8.90
pw = CW - 8.90
panel = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
    Inches(px), Inches(CY), Inches(pw), Inches(AH))
panel.fill.solid()
panel.fill.fore_color.rgb = BG_LIGHT
panel.line.fill.background()
ptb = slide.shapes.add_textbox(Inches(px + 0.15), Inches(CY + 0.15),
                                Inches(pw - 0.30), Inches(AH - 0.30))
ptf = ptb.text_frame
ptf.word_wrap = True
p = ptf.paragraphs[0]
r = p.add_run()
r.text = "Cloud specs"
r.font.size = Pt(16)
r.font.bold = True
r.font.color.rgb = DEEP_PURPLE
r.font.name = FONT

cloud_specs = [
    ("Backend:", "vLLM on H100"),
    ("Concurrency:", "10-20 simultaneous"),
    ("Cost target:", "< $10/run"),
    ("Provisioning:", "Ephemeral, auto-teardown"),
    ("CLI flag:", "--backend vllm"),
    ("Default:", "Ollama (offline)"),
]
for label_cs, val_cs in cloud_specs:
    p = ptf.add_paragraph()
    p.space_before = Pt(10)
    rl = p.add_run()
    rl.text = label_cs + " "
    rl.font.size = Pt(14)
    rl.font.bold = True
    rl.font.color.rgb = PURPLE
    rl.font.name = FONT
    rv = p.add_run()
    rv.text = val_cs
    rv.font.size = Pt(14)
    rv.font.color.rgb = TEXT_BODY
    rv.font.name = FONT

# Source note
src_tb = slide.shapes.add_textbox(Inches(ML + 0.10), Inches(BY - 0.35),
                                   Inches(8.50), Inches(0.30))
src_tf = src_tb.text_frame
p = src_tf.paragraphs[0]
r = p.add_run()
r.text = "339-opportunity Ford hierarchy catalog used as reference dataset"
r.font.size = Pt(14)
r.font.italic = True
r.font.color.rgb = TEXT_SUB
r.font.name = FONT


# ══════════════════════════════════════════════════════════════
# SLIDE 17 — What this means for your practice (Pattern O — bullet memo)
# ══════════════════════════════════════════════════════════════
slide = setup_content_slide(
    "Implications",
    "You do not need this exact tool. You need this approach.",
    ""
)
tb = slide.shapes.add_textbox(Inches(ML), Inches(CY), Inches(CW), Inches(AH))
tf = tb.text_frame
tf.word_wrap = True

memo_bullets = [
    ("The three-lens framework", " applies to any enterprise AI portfolio evaluation. "
     "Technical feasibility alone has never been the binding constraint on adoption."),
    ("Adoption-first weighting", " prevents the most common failure mode. "
     "When 45% of the score comes from decision density, financial gravity, "
     "impact proximity, and confidence signals, operationally dead features cannot score well."),
    ("Red flag triage", " catches data quality issues before expensive analysis. "
     "Five deterministic checks eliminate phantom opportunities and dead zones "
     "without consuming LLM inference time."),
    ("Grounded simulation", " prevents hallucinated specs from reaching implementation. "
     "A bundled knowledge base of real platform components means every generated "
     "artifact maps to something that actually exists."),
    ("Offline operation", " makes it suitable for air-gapped or sensitive client environments. "
     "The local Ollama path requires zero cloud connectivity. "
     "The cloud path is opt-in for speed, never required."),
]
for j, (bold_text, normal_text) in enumerate(memo_bullets):
    p = tf.add_paragraph() if j > 0 else tf.paragraphs[0]
    p.space_before = Pt(14) if j > 0 else Pt(0)
    rb = p.add_run()
    rb.text = f"  {bold_text}"
    rb.font.size = Pt(14)
    rb.font.bold = True
    rb.font.color.rgb = PURPLE
    rb.font.name = FONT
    rn = p.add_run()
    rn.text = normal_text
    rn.font.size = Pt(14)
    rn.font.color.rgb = TEXT_BODY
    rn.font.name = FONT


# ══════════════════════════════════════════════════════════════
# SLIDE 18 — Closing (Alt Cover — Layout 6)
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
for ph in slide.placeholders:
    idx = ph.placeholder_format.idx
    if idx == 0:
        ph.text = "Score what matters"
        for p in ph.text_frame.paragraphs:
            for r in p.runs:
                r.font.size = Pt(40)
                r.font.bold = True
                r.font.color.rgb = BLACK
                r.font.name = FONT
    elif idx == 13:
        ph.text = ("The question is not \"can we build it?\" "
                   "It is \"will they use it?\"\n\n"
                   "Adoption realism at 45% forces that question to the front "
                   "of every evaluation. The engine makes this repeatable, "
                   "defensible, and overnight-fast.\n\n"
                   "The best enterprise AI skill is the one that gets adopted.")
        for p in ph.text_frame.paragraphs:
            for r in p.runs:
                r.font.size = Pt(18)
                r.font.color.rgb = BLACK
                r.font.name = FONT


# ══════════════════════════════════════════════════════════════
# SAVE & VERIFY
# ══════════════════════════════════════════════════════════════
prs.save(OUTPUT)
cleanup_temp()
print(f"\nSaved: {OUTPUT}")
print(f"Slide count: {len(Presentation(OUTPUT).slides)}")

# Run verification
import subprocess
_py = sys.executable
verify_script = os.path.join(SKILL_DIR, "scripts", "verify_pptx.py")
brand_script = os.path.join(SKILL_DIR, "scripts", "brand_check.py")

print("\n=== Structural Verification ===")
subprocess.run([_py, verify_script, OUTPUT])
print("\n=== Brand Compliance ===")
subprocess.run([_py, brand_script, OUTPUT])
