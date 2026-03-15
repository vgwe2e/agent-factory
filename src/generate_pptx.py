#!/usr/bin/env python3
"""Generate Ford Top 10 Opportunities PowerPoint.

Design principle: the OPPORTUNITY is the hero. Show the problem, the flow,
the business case. Scores are reference data, not the headline.
"""

import csv
import json
import os
import re
import shutil
import textwrap
from collections import Counter

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
import matplotlib.ticker as ticker
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Colors ──
DARK_BLUE = RGBColor(0x00, 0x33, 0x66)
MED_BLUE = RGBColor(0x00, 0x5A, 0x9E)
ACCENT_GREEN = RGBColor(0x2E, 0x7D, 0x32)
ACCENT_ORANGE = RGBColor(0xE6, 0x7E, 0x22)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)
MED_GRAY = RGBColor(0xBB, 0xBB, 0xBB)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
PANEL_BG = RGBColor(0xF0, 0xF4, 0xFA)

# Flow diagram node colors (matplotlib)
FLOW_COLORS = {
    "Trigger": "#1565C0",
    "DDM/Crawlers": "#2E7D32",
    "Cortex Auto Forecast": "#6A1B9A",
    "AutoML": "#6A1B9A",
    "Cortex Monitoring": "#6A1B9A",
    "Remote Functions": "#E65100",
    "Action": "#00695C",
    "Outcome": "#F57F17",
}

def flow_color(label):
    for key, color in FLOW_COLORS.items():
        if key.lower() in label.lower():
            return color
    if "stream" in label.lower() or "ddm" in label.lower() or "crawl" in label.lower():
        return "#2E7D32"
    if "cortex" in label.lower() or "auto" in label.lower():
        return "#6A1B9A"
    if "pb" in label.lower() or "remote" in label.lower() or "process" in label.lower():
        return "#E65100"
    if "action" in label.lower():
        return "#00695C"
    if "outcome" in label.lower():
        return "#F57F17"
    if "trigger" in label.lower():
        return "#1565C0"
    return "#455A64"

BASE = os.path.dirname(os.path.abspath(__file__))
EVAL_DIR = os.path.join(BASE, "evaluation-vllm", "evaluation")
SIM_DIR = os.path.join(EVAL_DIR, "simulations")
PLANNING_DIR = os.path.join(BASE, "..", ".planning")

TOP10_SLUGS = [
    "demand-forecasting-analysis",
    "warehouse-inventory-management",
    "material-requirements-planning-mrp-integration",
    "multi-echelon-inventory-optimization",
    "production-planning-control",
    "inventory-positioning-optimization",
    "inbound-receiving-putaway",
    "cost-value-optimization-initiatives",
    "production-schedule-execution-monitoring",
    "quality-control-performance-testing",
]

TOP10_L3 = [
    "Demand Forecasting & Analysis",
    "Warehouse & Inventory Management",
    "Material Requirements Planning (MRP) Integration",
    "Multi-Echelon Inventory Optimization",
    "Production Planning & Control",
    "Inventory Positioning & Optimization",
    "Inbound Receiving & Putaway",
    "Cost & Value Optimization Initiatives",
    "Production Schedule Execution & Monitoring",
    "Quality Control & Performance Testing",
]

# ── Data loading ──

def load_tsv():
    with open(os.path.join(EVAL_DIR, "feasibility-scores.tsv")) as f:
        return list(csv.DictReader(f, delimiter="\t"))

def load_tier1_report():
    with open(os.path.join(EVAL_DIR, "tier1-report.md")) as f:
        return f.read()

def load_export():
    with open(os.path.join(PLANNING_DIR, "ford_hierarchy_v3_export.json")) as f:
        data = json.load(f)
    return {o["l3_name"]: o for o in data["project"]["l3_opportunities"]}

def load_mmd(slug):
    with open(os.path.join(SIM_DIR, slug, "decision-flow.mmd")) as f:
        return f.read()

def load_yaml(slug):
    with open(os.path.join(SIM_DIR, slug, "component-map.yaml")) as f:
        return f.read()

def get_top10_rows(all_rows):
    by_l3 = {r["l3_name"]: r for r in all_rows}
    return [by_l3[name] for name in TOP10_L3]

# ── Parsers ──

def parse_mermaid_flow(mmd_text):
    steps = []
    for line in mmd_text.strip().split("\n"):
        line = line.strip()
        m = re.match(r'A\[Trigger:\s*(.+?)\]', line)
        if m:
            steps.append(("Trigger", m.group(1).rstrip("].")))
            continue
        m = re.match(r'N\d+\[(?:Stream|Cortex|PB):\s*(.+?)\s*-\s*(.+?)\]', line)
        if m:
            steps.append((m.group(1).strip(), m.group(2).rstrip("].")))
            continue
        m = re.match(r'N\d+\{(?:PB):\s*(.+?)\s*-\s*(.+?)\}', line)
        if m:
            steps.append((m.group(1).strip(), m.group(2).rstrip("}.")))
            continue
        m = re.match(r'P\d+\[Action:\s*(.+?)\]', line)
        if m:
            steps.append(("Action", m.group(1)))
            continue
        m = re.match(r'Z\d+\[Outcome:\s*(.+?)\]', line)
        if m:
            steps.append(("Outcome", m.group(1)))
            continue
    if len(steps) > 7:
        steps = [steps[0]] + steps[1:6] + [steps[-1]]
    return steps

def parse_component_counts(yaml_text):
    counts = {}
    current = None
    for line in yaml_text.split("\n"):
        s = line.strip()
        if s.startswith("streams:"): current = "Streams"
        elif s.startswith("cortex:"): current = "Cortex"
        elif s.startswith("process_builder:"): current = "PB"
        elif s.startswith("agent_teams:"): current = "Agents"
        elif s.startswith("ui:"): current = "UI"
        elif s.startswith("- name:") and current:
            counts[current] = counts.get(current, 0) + 1
    return counts

# ── Flow diagram renderer ──

def render_flow_diagram(steps, output_path):
    """Render decision flow as a vertical box-and-arrow diagram."""
    n = len(steps)
    if n == 0:
        # Empty fallback
        fig, ax = plt.subplots(figsize=(5.5, 2))
        ax.text(0.5, 0.5, "No flow data", ha="center", va="center", fontsize=12, color="#999")
        ax.axis("off")
        fig.savefig(output_path, dpi=200, bbox_inches="tight", facecolor="white")
        plt.close(fig)
        return

    fig_height = max(3.0, n * 0.85 + 0.5)
    fig, ax = plt.subplots(figsize=(5.5, fig_height))
    ax.set_xlim(0, 10)
    ax.set_ylim(-n * 1.0 - 0.3, 0.5)
    ax.axis("off")

    box_width = 9.2
    box_height = 0.7
    x_center = 5.0
    x_left = x_center - box_width / 2

    for i, (label, desc) in enumerate(steps):
        y = -i * 1.0

        # Determine color
        color = flow_color(label)

        # Draw box
        rect = mpatches.FancyBboxPatch(
            (x_left, y - box_height / 2), box_width, box_height,
            boxstyle="round,pad=0.1",
            facecolor=color, edgecolor="white", linewidth=1.5, alpha=0.92
        )
        ax.add_patch(rect)

        # Label (bold, left side)
        label_text = label if len(label) <= 25 else label[:23] + ".."
        ax.text(x_left + 0.25, y + 0.08, label_text,
                fontsize=7.5, fontweight="bold", color="white", va="center",
                fontfamily="sans-serif")

        # Description (right side, wrapped)
        desc_wrapped = textwrap.shorten(desc, width=70, placeholder="...")
        ax.text(x_left + 0.25, y - 0.18, desc_wrapped,
                fontsize=6.5, color=(1, 1, 1, 0.85), va="center",
                fontfamily="sans-serif")

        # Arrow to next
        if i < n - 1:
            ax.annotate("", xy=(x_center, -(i + 1) * 1.0 + box_height / 2),
                        xytext=(x_center, y - box_height / 2),
                        arrowprops=dict(arrowstyle="-|>", color="#90A4AE", lw=2.0))

    plt.tight_layout(pad=0.1)
    fig.savefig(output_path, dpi=200, bbox_inches="tight", facecolor="white")
    plt.close(fig)

# ── Overview chart ──

def gen_overview_chart(top10_rows, export_data, path):
    """Grouped bar: composite + value for each opportunity."""
    fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(13, 5), gridspec_kw={"width_ratios": [1.2, 1]})

    names = []
    for r in top10_rows:
        n = r["l3_name"]
        for old, new in [("Optimization", "Opt."), ("Integration", "Integ."),
                         ("Requirements", "Req."), ("Forecasting", "Fcst."),
                         ("Management", "Mgmt."), ("Monitoring", "Mon."),
                         ("Performance", "Perf."), ("Initiatives", "Init."),
                         ("Execution", "Exec."), ("Planning", "Plan."),
                         ("Production", "Prod."), ("Positioning", "Pos.")]:
            n = n.replace(old, new)
        names.append(n[:28])

    # Left chart: 3-lens scores
    tech = [int(r["tech_total"]) for r in top10_rows]
    adopt = [int(r["adoption_total"]) for r in top10_rows]
    value = [int(r["value_total"]) for r in top10_rows]
    x = np.arange(len(names))
    w = 0.25

    ax1.bar(x - w, tech, w, label="Technical (/9)", color="#005A9E")
    ax1.bar(x, adopt, w, label="Adoption (/12)", color="#2E7D32")
    ax1.bar(x + w, value, w, label="Value (/6)", color="#E67E22")
    ax1.set_xticks(x)
    ax1.set_xticklabels(names, rotation=35, ha="right", fontsize=8)
    ax1.set_ylim(0, 14)
    ax1.yaxis.set_major_locator(ticker.MultipleLocator(3))
    ax1.legend(fontsize=8, loc="upper right")
    ax1.set_title("Three-Lens Scores", fontsize=12, fontweight="bold", color="#003366")
    for bars in ax1.containers:
        ax1.bar_label(bars, fontsize=6, padding=1)

    # Right chart: combined max value (horizontal bars)
    values_m = []
    for r in top10_rows:
        opp = export_data.get(r["l3_name"], {})
        v = opp.get("combined_max_value")
        values_m.append(v / 1e6 if v else 0)

    colors = ["#1565C0" if v > 50 else "#42A5F5" if v > 20 else "#90CAF9" for v in values_m]
    y = np.arange(len(names))
    ax2.barh(y, values_m, color=colors, height=0.6)
    ax2.set_yticks(y)
    ax2.set_yticklabels(names, fontsize=8)
    ax2.invert_yaxis()
    ax2.set_xlabel("$M", fontsize=10)
    ax2.set_title("Combined Max Value ($M)", fontsize=12, fontweight="bold", color="#003366")
    for i, v in enumerate(values_m):
        if v > 0:
            ax2.text(v + 1, i, f"${v:.0f}M", va="center", fontsize=8, fontweight="bold")

    plt.tight_layout()
    fig.savefig(path, dpi=250, bbox_inches="tight", facecolor="white")
    plt.close(fig)

# ── Slide helpers ──

def set_bg(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE

def header_bar(slide, text, height=Inches(0.65)):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.4)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = WHITE

def tbox(slide, left, top, width, height, text, size=12, bold=False, color=DARK_GRAY, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = "Calibri"
    p.alignment = align
    return tb

def mlines(slide, left, top, width, height, lines, size=11, color=DARK_GRAY):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = Pt(2)

def add_info_table(slide, left, top, width, rows_data):
    """Add a 2-column info table (label | value). Returns bottom y position."""
    n = len(rows_data)
    table_shape = slide.shapes.add_table(n, 2, left, top, width, Inches(n * 0.38))
    table = table_shape.table
    table.columns[0].width = Inches(1.8)
    table.columns[1].width = width - Inches(1.8)

    for i, (label, value) in enumerate(rows_data):
        # Label cell
        cell_l = table.cell(i, 0)
        cell_l.text = label
        cell_l.fill.solid()
        cell_l.fill.fore_color.rgb = RGBColor(0xE3, 0xEA, 0xF5)
        for p in cell_l.text_frame.paragraphs:
            p.font.size = Pt(9)
            p.font.bold = True
            p.font.color.rgb = DARK_BLUE
            p.font.name = "Calibri"

        # Value cell
        cell_v = table.cell(i, 1)
        cell_v.text = value
        cell_v.fill.solid()
        cell_v.fill.fore_color.rgb = WHITE
        for p in cell_v.text_frame.paragraphs:
            p.font.size = Pt(9)
            p.font.color.rgb = DARK_GRAY
            p.font.name = "Calibri"

    return top + Inches(n * 0.38)

def add_scores_strip(slide, left, top, row):
    """Compact 1-row scores table at bottom."""
    tech_denominator = "9" if "data_readiness" in row or "archetype_conf" in row else "3"
    cols = ["Tech", "Adopt", "Value", "Composite"]
    vals = [
        f"{row['tech_total']}/{tech_denominator}", f"{row['adoption_total']}/12", f"{row['value_total']}/6",
        row["composite"],
    ]

    if "data_readiness" in row:
        cols.append("DR")
        vals.append(row.get("data_readiness", ""))

    cols.append("PF")
    vals.append(row.get("platform_fit", ""))

    if "archetype_conf" in row:
        cols.append("AC")
        vals.append(row.get("archetype_conf", ""))

    cols.extend(["DD", "FG", "IP", "CS", "VD", "SV"])
    vals.extend([
        row.get("decision_density", ""),
        row.get("financial_gravity", ""),
        row.get("impact_proximity", ""),
        row.get("confidence_signal", ""),
        row.get("value_density", ""),
        row.get("simulation_viability", ""),
    ])

    table_shape = slide.shapes.add_table(2, len(cols), left, top, Inches(12.5), Inches(0.55))
    table = table_shape.table

    for i in range(len(cols)):
        table.columns[i].width = Inches(12.5 / len(cols))

    for i, (col, val) in enumerate(zip(cols, vals)):
        # Header
        cell_h = table.cell(0, i)
        cell_h.text = col
        cell_h.fill.solid()
        cell_h.fill.fore_color.rgb = RGBColor(0x90, 0xA4, 0xAE)
        for p in cell_h.text_frame.paragraphs:
            p.font.size = Pt(7)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.alignment = PP_ALIGN.CENTER

        # Value
        cell_v = table.cell(1, i)
        cell_v.text = str(val)
        cell_v.fill.solid()
        cell_v.fill.fore_color.rgb = RGBColor(0xEC, 0xEF, 0xF1)
        for p in cell_v.text_frame.paragraphs:
            p.font.size = Pt(8)
            p.font.color.rgb = DARK_GRAY
            p.alignment = PP_ALIGN.CENTER
            if i == 3:  # Composite
                p.font.bold = True
                p.font.color.rgb = DARK_BLUE

# ── Slide builders ──

def slide_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)

    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(3.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()

    tbox(slide, Inches(1), Inches(0.8), Inches(11), Inches(1), "Ford Motor Company", size=42, bold=True, color=WHITE)
    tbox(slide, Inches(1), Inches(1.7), Inches(11), Inches(0.8), "Top 10 AI Skill Opportunities",
         size=32, color=RGBColor(0xA0, 0xC4, 0xE8))
    tbox(slide, Inches(1), Inches(2.5), Inches(11), Inches(0.6), "Aera Skill Feasibility Assessment",
         size=20, color=RGBColor(0xCC, 0xDD, 0xEE))

    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(3.5), Inches(13.333), Inches(0.7))
    bar.fill.solid()
    bar.fill.fore_color.rgb = MED_BLUE
    bar.line.fill.background()
    tbox(slide, Inches(1), Inches(3.55), Inches(11), Inches(0.5),
         "336 Evaluated  |  178 Promoted  |  146 Simulated  |  ~$2.4B Combined Value",
         size=18, bold=True, color=WHITE)

    tbox(slide, Inches(1), Inches(4.8), Inches(11), Inches(0.5), "March 2026", size=16, color=DARK_GRAY)
    tbox(slide, Inches(1), Inches(5.3), Inches(11), Inches(0.4),
         "Aera Skill Feasibility Engine v1.2  |  vLLM on RunPod H100",
         size=11, color=MED_GRAY)

def slide_exec_summary(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    header_bar(slide, "Executive Summary")

    tbox(slide, Inches(0.5), Inches(0.9), Inches(6), Inches(0.4),
         "Three-Lens Scoring", size=18, bold=True, color=DARK_BLUE)
    mlines(slide, Inches(0.5), Inches(1.4), Inches(6), Inches(2.5), [
        "Technical Feasibility (30%) -- Data readiness, platform fit, archetype confidence",
        "Adoption Realism (45%) -- Decision density, financial gravity, impact proximity, confidence",
        "Value & Efficiency (25%) -- Value density, simulation viability",
        "",
        "Composite = (tech/9 x 0.30) + (adoption/12 x 0.45) + (value/6 x 0.25)",
        "Simulation threshold: composite >= 0.60",
    ], size=11)

    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7), Inches(0.9), Inches(5.8), Inches(3.0))
    shape.fill.solid()
    shape.fill.fore_color.rgb = PANEL_BG
    shape.line.color.rgb = MED_BLUE
    shape.line.width = Pt(1)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.3)
    tf.margin_top = Inches(0.2)
    p = tf.paragraphs[0]
    p.text = "Pipeline Results"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    for s in ["336 total opportunities evaluated across Ford hierarchy",
              "178 promoted to simulation (53%)",
              "146 simulations completed with full artifact generation",
              "Top 10: all >= 0.82 composite, all HIGH confidence",
              "Combined max value of top 10: ~$2.4B",
              "Tier 1: 39 | Tier 2: 103 | Tier 3: 220",
              "DETERMINISTIC: 57% | AGENTIC: 43%"]:
        p2 = tf.add_paragraph()
        p2.text = "  " + s
        p2.font.size = Pt(11)
        p2.font.color.rgb = DARK_GRAY
        p2.space_before = Pt(3)

def slide_leaderboard(prs, top10, export_data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    header_bar(slide, "Top 10 Opportunities")

    cols = 7
    table_shape = slide.shapes.add_table(11, cols, Inches(0.2), Inches(0.9), Inches(12.9), Inches(6.0))
    table = table_shape.table
    widths = [Inches(0.5), Inches(3.5), Inches(1.6), Inches(1.2), Inches(1.3), Inches(1.5), Inches(3.3)]
    for i, w in enumerate(widths):
        table.columns[i].width = w

    headers = ["#", "Skill / Use Case", "Domain", "Score", "Max Value", "Complexity", "Problem (excerpt)"]
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = DARK_BLUE
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(9)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.alignment = PP_ALIGN.CENTER

    score_bg = {0.88: RGBColor(0xC6, 0xEF, 0xCE), 0.84: RGBColor(0xDA, 0xF2, 0xDA), 0.82: RGBColor(0xE8, 0xF5, 0xE9)}

    for idx, row in enumerate(top10):
        r = idx + 1
        opp = export_data.get(row["l3_name"], {})
        composite = float(row["composite"])
        bg = score_bg.get(composite, RGBColor(0xF5, 0xF5, 0xF5))
        max_val = opp.get("combined_max_value")
        val_str = f"${max_val/1e6:.0f}M" if max_val else "N/A"
        summary = opp.get("opportunity_summary", "")
        excerpt = summary[:120] + "..." if len(summary) > 120 else summary

        vals = [str(r), row["opportunity_name"], row["l1_name"],
                f"{composite:.2f}", val_str, opp.get("implementation_complexity", ""),
                excerpt]
        for c, v in enumerate(vals):
            cell = table.cell(r, c)
            cell.text = v
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(8)
                p.font.color.rgb = DARK_GRAY
                if c == 3:
                    p.font.bold = True
                    p.font.color.rgb = DARK_BLUE
                if c in (0, 3, 4, 5):
                    p.alignment = PP_ALIGN.CENTER

def slide_overview_chart(prs, top10, export_data, chart_path):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    header_bar(slide, "Scoring & Value Overview")
    gen_overview_chart(top10, export_data, chart_path)
    slide.shapes.add_picture(chart_path, Inches(0.15), Inches(0.8), Inches(13.0), Inches(6.4))

def slide_opportunity(prs, rank, row, slug, tier1_text, export_data, tmp_dir):
    """One dense slide: LEFT = info table, RIGHT = flow diagram, BOTTOM = scores strip."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)

    l3 = row["l3_name"]
    opp = export_data.get(l3, {})
    composite = float(row["composite"])

    # ── Header bar with key stats ──
    max_val = opp.get("combined_max_value")
    val_str = f"${max_val/1e6:.0f}M" if max_val else "N/A"
    quick_win = opp.get("quick_win", False)
    qw_tag = "  |  QUICK WIN" if quick_win else ""

    header_text = f"#{rank}  {row['opportunity_name']}    {composite:.2f}  |  {val_str}  |  {row['archetype']}{qw_tag}"
    header_bar(slide, header_text, height=Inches(0.6))

    # Hierarchy path
    tbox(slide, Inches(0.4), Inches(0.65), Inches(12), Inches(0.22),
         f"{row['l1_name']}  >  {row['l2_name']}  >  {l3}",
         size=9, bold=True, color=MED_BLUE)

    # ── LEFT: Info table (full problem statement + business case) ──
    summary = opp.get("opportunity_summary", "No summary available.")
    rationale = opp.get("rationale", "No rationale available.")
    complexity = opp.get("implementation_complexity", "N/A")
    l4_count = opp.get("l4_count", "?")
    hv_l4 = opp.get("high_value_l4_count", "?")
    support_arch = opp.get("supporting_archetypes", [])
    arch_str = f"{opp.get('lead_archetype', row['archetype'])}"
    if support_arch:
        arch_str += f" + {', '.join(support_arch)}"

    # Component counts
    yaml_text = load_yaml(slug)
    counts = parse_component_counts(yaml_text)
    comp_str = "  |  ".join(f"{k}: {v}" for k, v in counts.items() if v > 0) or "N/A"

    table_rows = [
        ("Skill Name", row["opportunity_name"]),
        ("Problem Statement", summary),
        ("Value Rationale", rationale),
        ("Combined Max Value", val_str),
        ("Complexity", complexity),
        ("Quick Win", "Yes" if quick_win else "No"),
        ("L4 Activities", f"{l4_count} total, {hv_l4} high-value"),
        ("Archetypes", arch_str),
        ("Aera Components", comp_str),
    ]

    add_info_table(slide, Inches(0.3), Inches(0.92), Inches(6.8), table_rows)

    # ── RIGHT: Flow diagram ──
    mmd_text = load_mmd(slug)
    flow_steps = parse_mermaid_flow(mmd_text)

    flow_path = os.path.join(tmp_dir, f"flow_{rank}.png")
    render_flow_diagram(flow_steps, flow_path)

    tbox(slide, Inches(7.3), Inches(0.92), Inches(5.5), Inches(0.22),
         "Decision Flow", size=11, bold=True, color=DARK_BLUE)
    slide.shapes.add_picture(flow_path, Inches(7.3), Inches(1.2), Inches(5.7), Inches(5.2))

    # ── BOTTOM: Scores strip ──
    add_scores_strip(slide, Inches(0.3), Inches(6.85), row)

def slide_next_steps(prs, top10, export_data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    header_bar(slide, "Recommended Next Steps")

    wave1 = [(r, export_data.get(r["l3_name"], {})) for r in top10 if float(r["composite"]) >= 0.88]
    wave2 = [(r, export_data.get(r["l3_name"], {})) for r in top10 if 0.83 <= float(r["composite"]) < 0.88]
    wave3 = [(r, export_data.get(r["l3_name"], {})) for r in top10 if float(r["composite"]) < 0.83]

    def fmt_wave(items):
        lines = []
        for r, o in items:
            v = o.get("combined_max_value")
            vs = f"  ${v/1e6:.0f}M" if v else ""
            qw = "  [QUICK WIN]" if o.get("quick_win") else ""
            cx = o.get("implementation_complexity", "")
            lines.append(f"   {r['l3_name']}  ({r['l1_name']}){vs}  {cx}{qw}")
        return lines

    lines = ["Wave 1 -- Highest Confidence (0.88):"] + fmt_wave(wave1)
    lines += ["", "Wave 2 -- Strong Candidates (0.84):"] + fmt_wave(wave2)
    lines += ["", "Wave 3 -- Solid Foundation (0.82):"] + fmt_wave(wave3)

    mlines(slide, Inches(0.5), Inches(0.9), Inches(6.5), Inches(6), lines, size=11)

    # Domain clusters
    tbox(slide, Inches(7.2), Inches(0.9), Inches(5.5), Inches(0.3),
         "Domain Clusters", size=16, bold=True, color=DARK_BLUE)
    l1_counts = Counter(r["l1_name"] for r in top10)
    dep = []
    for l1, cnt in l1_counts.most_common():
        if cnt > 1:
            related = [r["l3_name"] for r in top10 if r["l1_name"] == l1]
            dep.append(f"{l1} ({cnt} opportunities):")
            for name in related:
                dep.append(f"   {name}")
            dep.append("")
    if dep:
        mlines(slide, Inches(7.2), Inches(1.3), Inches(5.5), Inches(5), dep, size=10)

# ── Main ──

def main():
    print("Loading data...")
    all_rows = load_tsv()
    top10 = get_top10_rows(all_rows)
    tier1_text = load_tier1_report()
    export_data = load_export()

    tmp_dir = os.path.join(EVAL_DIR, "_tmp_charts")
    os.makedirs(tmp_dir, exist_ok=True)

    print("Building presentation...")
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide_title(prs)
    print("  [1] Title")

    slide_exec_summary(prs)
    print("  [2] Exec Summary")

    slide_leaderboard(prs, top10, export_data)
    print("  [3] Leaderboard")

    chart_path = os.path.join(tmp_dir, "overview.png")
    slide_overview_chart(prs, top10, export_data, chart_path)
    print("  [4] Scoring & Value Chart")

    for i, (row, slug) in enumerate(zip(top10, TOP10_SLUGS)):
        rank = i + 1
        slide_opportunity(prs, rank, row, slug, tier1_text, export_data, tmp_dir)
        print(f"  [{4+rank}] #{rank} {row['l3_name']}")

    slide_next_steps(prs, top10, export_data)
    print("  [15] Next Steps")

    output = os.path.join(EVAL_DIR, "ford-top10-opportunities.pptx")
    prs.save(output)
    shutil.rmtree(tmp_dir, ignore_errors=True)

    p = Presentation(output)
    sz = os.path.getsize(output)
    print(f"\nSaved: {output}")
    print(f"  {len(p.slides)} slides, {sz:,} bytes, {sum(1 for s in p.slides for sh in s.shapes if hasattr(sh, 'image'))} images")

if __name__ == "__main__":
    main()
